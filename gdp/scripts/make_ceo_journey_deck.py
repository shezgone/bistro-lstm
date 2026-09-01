# -*- coding: utf-8 -*-
"""대표 보고 덱 — 한국은행 협업 전체 여정 (10장, 실무자 스타일, 통계학 배경 청중)"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

F = "Apple SD Gothic Neo"
INK = RGBColor(0x26,0x26,0x26); GREY = RGBColor(0x6B,0x6B,0x6B)
LINE = RGBColor(0xC9,0xC9,0xC9); NAVY = RGBColor(0x24,0x34,0x47)
GREEN = RGBColor(0x00,0x8A,0x3E); BGGREY = RGBColor(0xF5,0xF5,0xF4)
WHITE = RGBColor(0xFF,0xFF,0xFF); HL = RGBColor(0xEE,0xF5,0xF0)
WARM = RGBColor(0xB0,0x53,0x2F); BLUE = RGBColor(0x1C,0x5C,0xAB)

SC = "/private/tmp/claude-502/-Users-user-vibe-bistro-lstm/031ca4b9-003d-4cb2-afaa-06353432dc9b/scratchpad"
prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
PAGE = [0]
def new_slide():
    PAGE[0] += 1
    return prs.slides.add_slide(prs.slide_layouts[6])

def runs(s, x, y, w, h, lines, size=10, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp=2):
    tf = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.02); tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01); tf.margin_bottom = Inches(0.01)
    for i, rs in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp)
        for t, b, c, z in rs:
            r = p.add_run(); r.text = t
            r.font.name = F; r.font.size = Pt(z if z else size)
            r.font.bold = b; r.font.color.rgb = c if c else INK
            rp = r._r.get_or_add_rPr(); rp.append(rp.makeelement(qn('a:ea'), {'typeface': F}))
    return tf

def hline(s, x, y, w, color=LINE, pt=0.75):
    ln = s.shapes.add_connector(1, Inches(x), Inches(y), Inches(x+w), Inches(y))
    ln.line.color.rgb = color; ln.line.width = Pt(pt); return ln

def rect(s, x, y, w, h, fill=None, line=None):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False; return shp

def header(s, tag, msg):
    runs(s, 0.6, 0.32, 11.5, 0.3, [[(tag, True, GREY, 10.5)]])
    runs(s, 0.6, 0.6, 12.1, 0.6, [[(msg, True, NAVY, 17)]])
    hline(s, 0.6, 1.22, 12.13, NAVY, 1.2)

def footnote(s, text):
    runs(s, 0.6, 7.1, 11.6, 0.3, [[("주: " + text, False, GREY, 8)]])
    runs(s, 12.85, 7.1, 0.4, 0.25, [[(str(PAGE[0]), False, GREY, 9)]], align=PP_ALIGN.RIGHT)

def sidebar(s, x, y, h, color=NAVY):
    rect(s, x, y, 0.045, h, color, None)

def notes(s, text):
    s.notes_slide.notes_text_frame.text = text

RULE = "공통 규약: 한은 실시간 빈티지 · 속보치 기준 · 전망주차 w[-19,-1] 평균 RMSE · 2018Q1~2025Q4(32개 분기). 낮을수록 정확."

# ============ 1. 표지 ============
s = new_slide()
runs(s, 0.75, 0.55, 8.0, 0.35, [[("네이버클라우드 AX Forward Lab", True, GREY, 12)]])
runs(s, 0.75, 2.5, 11.8, 1.7, [
    [("한국은행 GDP Nowcasting 협업 — 3개월의 여정", True, NAVY, 31)],
    [("", False, INK, 8)],
    [("국면전환 가설에서 파운데이션 모델 적응학습까지, 그리고 공동 연구 어젠다", False, INK, 16)],
])
hline(s, 0.78, 4.35, 4.2, NAVY, 1.6)
runs(s, 0.78, 4.55, 11.2, 1.0, [
    [("40여 개 모형 구성 · 실험 스크립트 50여 종 · 동일 규약 실시간 재현 채점 608발 × 32개 분기.", False, GREY, 12)],
    [("정확도 신기록(-2.3%)과 함께, 한은이 협업 범위 확대를 먼저 제안해온 상태입니다.", False, GREY, 12)],
], sp=3)
runs(s, 0.78, 6.6, 8.0, 0.35, [[("2026. 9.  |  대표 보고  |  작성: 김용민", False, GREY, 11)]])
notes(s, "3개월 여정의 한 줄 요약: 정면 승부(국면전환) → 통계적 벽 확인 → 축 전환(조건부 가치) → 고객이 협업 확대를 제안하는 위치까지. 발표 시간 10분 기준, 2·5·9번 장이 핵심.")

# ============ 2. 과제와 여정 지도 ============
s = new_slide()
header(s, "1  |  과제와 여정", "요약 — 정확도 경쟁에서 출발해, 평가 기준을 함께 다시 세우는 데까지 왔습니다")
phases = [
    ("6월", "착수·재현", "한은 측의 시계열 파운데이션 모델(BISTRO, BIS WP 1337) 해석 질의에서 출발. FDE Consulting & Planning팀(김지혜·김민승)과 초기 방향성을 함께 수립하고, 현행 시스템(DFM+XGBoost)을 실시간 빈티지로 재현해 검증 기반을 세움", BLUE),
    ("7월 초", "국면전환 연구", "국면별 모델 교대(v1→v3)로 최고 성능 도달 — 그러나 방법론 수용성 문제로 채택 보류", WARM),
    ("7월 말", "축 전환", "\"32개 분기에선 누구도 우승을 증명 못 한다\"(검정력 분석) → 조건부 가치·분포로 성과 재정의", NAVY),
    ("8월 초", "기준선 격변", "한은이 자체 코드의 단위 결함을 개정 — 전 실험 즉시 재검증, XGBoost 단독이 새 기준선", GREY),
    ("8월 말", "적응학습 실증", "한은 제안(적응학습)을 당사가 첫 실증 — LoRA + 임계 규칙으로 32분기 신기록(0.733)", GREEN),
]
y = 1.5
for when, title, desc, c in phases:
    sidebar(s, 0.6, y, 0.88, c)
    runs(s, 0.85, y, 1.35, 0.4, [[(when, True, c, 12)]])
    runs(s, 2.25, y, 2.2, 0.4, [[(title, True, NAVY, 12)]])
    runs(s, 4.55, y, 8.1, 0.8, [[(desc, False, INK, 10.5)]])
    y += 1.06
footnote(s, RULE)
notes(s, "여정의 구조가 메시지: 우리는 실패를 숨기지 않고 기록했고, 그 기록(무엇이 왜 안 되는가)이 지금 한은과의 공동 논문 재료가 됨. 8월 기준선 격변 때 당일 재검증으로 대응한 것이 신뢰의 전환점. (구두 보충용 — 문서 명시 금지: BISTRO 논문의 1저자가 현 신임 총재(BIS 재직 시)이며, 디지털혁신실의 해석 질의가 이 프로젝트의 기점. 8월의 원본 실측·LoRA 개선으로 출발 질문에 답을 완성한 수미상관 구조.)")

# ============ 3. 검증 체계 ============
s = new_slide()
header(s, "2  |  검증 체계", "모든 주장의 바닥에 하나의 채점 규약 — 실시간 재현, 예외 없음")
runs(s, 0.6, 1.45, 6.0, 4.6, [
    [("설계 원칙", True, NAVY, 12.5)],
    [("· 실시간 빈티지: 각 예측 시점에 '그때 존재했던' 데이터만 사용", False, INK, 10.5)],
    [("  — look-ahead 원천 차단, 개정 이력 보존", False, GREY, 9.5)],
    [("· 정답 = 속보치, 분기당 19회 주간 갱신 예측 (총 608발)", False, INK, 10.5)],
    [("· 2018Q1~2025Q4, 코로나 급변기 포함 전 구간", False, INK, 10.5)],
    [("", False, INK, 6)],
    [("통계적 안전장치", True, NAVY, 12.5)],
    [("· 검정: Diebold-Mariano(HAC), Model Confidence Set", False, INK, 10.5)],
    [("· 검정력 분석: 이 표본에선 6% 개선도 p≈0.09 — 사전 확인", False, INK, 10.5)],
    [("· 사전 등록: 임계값·설정을 실험 전 문서로 고정", False, INK, 10.5)],
    [("· seed 3회 반복 · 선택편의 명시 · 경계 민감도 검사", False, INK, 10.5)],
    [("· 단일 실행/단일 분기 수치는 판정에 사용하지 않음 (원칙)", False, INK, 10.5)],
])
rect(s, 7.0, 1.45, 5.75, 4.4, BGGREY, None)
sidebar(s, 7.0, 1.45, 4.4, NAVY)
runs(s, 7.25, 1.65, 5.3, 4.1, [
    [("이 규약이 실전에서 일한 순간들", True, NAVY, 12)],
    [("", False, INK, 4)],
    [("· 한은 코드 개정(8월) 때: 규약이 이미 서 있어 당일로", False, INK, 10.5)],
    [("  40여 구성 전면 재채점 — 기준선 교체에 즉시 대응", False, INK, 10.5)],
    [("· seed 미전파 결함 2건을 반복 원칙이 적발 — 행운의", False, INK, 10.5)],
    [("  단일 실행이 '신기록'으로 보고될 뻔한 것을 차단", False, INK, 10.5)],
    [("· 한은의 자체 MCS 검정이 당사 검정력 분석과 동일", False, INK, 10.5)],
    [("  결론에 도달 — 양측이 같은 언어로 논의하게 됨", False, INK, 10.5)],
], sp=3)
footnote(s, RULE)
notes(s, "통계학 배경 청중 핵심 장. 메시지: 점수보다 채점의 무결성에 먼저 투자했고, 그것이 고객 신뢰(코드 개정 당일 대응)와 연구 신뢰(seed 사고 적발) 양쪽에서 배당을 지급함.")

# ============ 4. 1막: 국면전환 ============
s = new_slide()
header(s, "3  |  1막 · 국면전환 연구", "성능은 최고였으나, 채택되지 않는 이유를 배웠습니다")
runs(s, 0.6, 1.45, 6.1, 4.4, [
    [("접근 — 국면전환(regime-switching) 계보", True, NAVY, 12.5)],
    [("경기 국면(급변·안정·반등)마다 최적 모형이 다르다는 가설로", False, INK, 10.5)],
    [("국면 감지 + 모델 교대 구조를 3차례 고도화 (v1→v3) —", False, INK, 10.5)],
    [("잠재 국면이 관측을 지배한다는 계량경제학의 고전적 발상.", False, INK, 10.5)],
    [("", False, INK, 5)],
    [("결과", True, NAVY, 12.5)],
    [("· v3 = 0.718 — 전체 여정 통틀어 최고 성능 (방향은 옳았음)", False, INK, 10.5)],
    [("· DM p=0.088 — 표본 한계로 유의성 문턱 미달", False, INK, 10.5)],
    [("", False, INK, 5)],
    [("채택 보류의 이유 (성능이 아니라 수용성)", True, WARM, 12.5)],
    [("한은 실무 책임자: \"현재 시점에 국면을 실시간으로 단정하는", False, INK, 10.5)],
    [("것이 맞는가\" — 중앙은행 운영 요건(설명책임)과의 충돌.", False, INK, 10.5)],
    [("증명을 계속하는 대신 자산화(부품 재사용·논문화)로 전환.", False, INK, 10.5)],
])
rect(s, 7.1, 1.45, 5.65, 4.4, HL, None)
sidebar(s, 7.1, 1.45, 4.4, GREEN)
runs(s, 7.35, 1.65, 5.2, 4.1, [
    [("이 연구가 남긴 자산 (이후 전부 재사용됨)", True, GREEN, 12)],
    [("", False, INK, 4)],
    [("· \"반등 국면에서 사전학습 모델이 강하다\"는 발견", False, INK, 10.5)],
    [("  → 이후 FM 검증·조기 슬롯 설계의 출발점", False, INK, 10.5)],
    [("· 국면 신호(심리 저점 통과 등) → 설명 리포트 재료", False, INK, 10.5)],
    [("· '이산 국면 단정' 대신 '전망시계 조건화'라는", False, INK, 10.5)],
    [("  수용 가능한 프레임 발견 — 최종 구성의 뼈대", False, INK, 10.5)],
    [("", False, INK, 4)],
    [("잠재 국면을 명시적으로 감지하는 대신, 국면 정보가", False, INK, 10.5)],
    [("자연히 실리는 조건(전망 주차·일별 신호)으로 우회 —", False, INK, 10.5)],
    [("같은 직관의 운영 가능한 번역이었습니다.", True, INK, 10.5)],
], sp=3)
footnote(s, RULE + " v3 상세: 반등 국면 감지 시 사전학습 모델 결합으로 교대하는 3-arm 구조.")
notes(s, "(발표자 참고 — 본문 언급 없음) 대표님이 과거 한은 맥락에서 HMM을 언급한 바 있음. 이 국면전환 연구가 정확히 그 계보(Markov-switching)이며 성능상 전 여정 최고(0.718)였음 — 질문이 나오면 구두로 '대표님이 보셨던 방향이 성능으로는 옳았고, 수용성 문제로 전망시계 조건화로 번역해 최종 구성에 계승했다'로 답변. 선제 언급은 하지 않음.")

# ============ 5. 2막: 검정력의 벽 ============
s = new_slide()
header(s, "4  |  2막 · 통계적 벽", "\"이 표본에서는 어떤 모형의 우위도 증명이 어렵다\" — 그래서 평가의 축을 바꿨습니다")
runs(s, 0.6, 1.45, 6.1, 4.4, [
    [("검정력 분석이 말해준 것", True, NAVY, 12.5)],
    [("· 표본 = 32개 분기. 주간 19발은 분기 내 강한 상관 —", False, INK, 10.5)],
    [("  유효 표본은 사실상 32", False, INK, 10.5)],
    [("· 이 표본에서 6.2% 개선도 DM p=0.088 — 5% 문턱을", False, INK, 10.5)],
    [("  넘으려면 표본이 수 배 필요", False, INK, 10.5)],
    [("· 한은의 자체 MCS: 비교 12개 모형 전원이 5% 신뢰집합에", False, INK, 10.5)],
    [("  잔류 — 독립적으로 같은 결론", False, INK, 10.5)],
    [("", False, INK, 5)],
    [("따라서 — 점 RMSE '순위 경쟁'은 판정 불가능한 게임", True, WARM, 12)],
])
runs(s, 7.1, 1.45, 5.65, 0.4, [[("성과 축의 재정의 (7월 말 결정)", True, NAVY, 12.5)]])
hline(s, 7.1, 1.85, 5.65, LINE, 1.0)
runs(s, 7.1, 2.0, 5.65, 3.8, [
    [("① 점추정 — ", True, INK, 11), ("'비열등 + 개선 방향' 수위로만 주장", False, INK, 10.5)],
    [("② 조건부 가치 — ", True, INK, 11), ("\"언제 무엇이 기여하는가\"는 유의성", False, INK, 10.5)],
    [("   없이도 성립하는 기술적(descriptive) 발견", False, INK, 10.5)],
    [("③ 분포 — ", True, INK, 11), ("예측구간 커버리지는 즉시 검증 가능한 축", False, INK, 10.5)],
    [("   (컨포멀 보정으로 명목 80% → 실측 81~87% 달성)", False, INK, 10.5)],
    [("④ 방법론 — ", True, INK, 11), ("검정력·선택편의·재현성 자체를 기여로", False, INK, 10.5)],
    [("", False, INK, 5)],
    [("이 재정의가 이후 모든 산출물의 문법이 됐고,", False, INK, 10.5)],
    [("공동 논문 프레임(\"언제 가치를 더하는가\")의 원형입니다.", True, INK, 10.5)],
], sp=3)
footnote(s, RULE)
notes(s, "통계학자 청중에게 가장 공명할 장. '못 이겨서 물러선 것'이 아니라 '검정력 계산이 먼저 있었고, 판정 가능한 질문으로 옮긴 것'. 한은 MCS가 독립적으로 같은 결론에 도달한 것이 이 판단의 외부 검증.")

# ============ 6. 3막: 기준선 격변 ============
s = new_slide()
header(s, "5  |  3막 · 기준선 격변", "한은의 코드 개정에 당일 재검증으로 대응했습니다")
runs(s, 0.6, 1.45, 6.0, 4.4, [
    [("무슨 일이 있었나 (8월 초)", True, NAVY, 12.5)],
    [("· 한은이 자체 ML 예측치의 단위 정합성 결함을 발견·개정", False, INK, 10.5)],
    [("  — 표준화값이 원단위와 혼합 채점되던 문제", False, INK, 10.5)],
    [("· 종전 공식 기준 'DFM+XGBoost 0.765'가 무효화,", False, INK, 10.5)],
    [("  당사의 종전 기록·게이트 수치도 함께 무효", False, INK, 10.5)],
    [("", False, INK, 5)],
    [("당사 대응", True, NAVY, 12.5)],
    [("· 개정 코드 수신 당일 40여 구성 전면 재채점", False, INK, 10.5)],
    [("· 신 기준선 확정: XGBoost 단독 0.750 —", False, INK, 10.5)],
    [("  \"결합이 항상 낫다\"던 종전 서사가 결함의 산물이었음을 확인", False, INK, 10.5)],
    [("· 한은의 대규모 후보 탐색(911개)에는 선택편의 관점을", False, INK, 10.5)],
    [("  정중히 제기 — 그들 자신의 MCS 논리로", False, INK, 10.5)],
])
s.shapes.add_picture(f"{SC}/bokf_leader.png", Inches(6.85), Inches(1.6), width=Inches(6.1))
footnote(s, RULE + " 개정 전 수치와의 혼용 없음 — 이후 모든 수치는 개정(schema v2) 기준.")
notes(s, "FDE 가치가 가장 선명한 장. 고객 결함을 비난 없이 '개정'으로 프레임하고, 우리 무효 수치(0.746 신기록 포함)도 함께 폐기하는 대칭성을 보임 — 이 대응 이후 한은이 자료를 자기 레포에 커밋하고 협업 확대를 제안하는 관계로 전환됨.")

# ============ 7. 4막: FM 전수 검증 ============
s = new_slide()
header(s, "6  |  4막 · 파운데이션 모델 전수 검증", "딥러닝은 조건부로 유효했습니다 — BIS 공개 모델 원본 실측 포함")
runs(s, 0.6, 1.45, 6.0, 4.5, [
    [("전수 검증의 결론", True, NAVY, 12.5)],
    [("· 직접학습 신경망(LSTM·NCDE 등): 전멸 — 월 140행 소표본", False, INK, 10.5)],
    [("· 사전학습 FM(Chronos-2·Moirai/BISTRO): 유효 —", False, INK, 10.5)],
    [("  단 '어느 조건에서'가 갈림 (조기 주차·반등기 한정)", False, INK, 10.5)],
    [("· BIS 공개 거시 FM(BISTRO, WP1337)을 확보해 동일 규약", False, INK, 10.5)],
    [("  직접 실측 — 한은 관심 질문에 원본 기준으로 답", False, INK, 10.5)],
    [("", False, INK, 5)],
    [("핵심 진단 — '힌트 흡수력'", True, NAVY, 12.5)],
    [("같은 일별 신호(주가·환율)를 줘도 활용도가 다름:", False, INK, 10.5)],
    [("Chronos-2 -3.4% vs BISTRO 무이득 (절제 실험)", True, INK, 10.5)],
    [("원인 = 사전학습 채점 범위 차이 (타깃만 vs 전 변량)", False, INK, 10.5)],
    [("— 주차별 오차 평탄성(한은 진단)의 원인을 규명", False, INK, 10.5)],
])
s.shapes.add_picture(f"{SC}/weekly_absorb.png", Inches(6.8), Inches(1.7), width=Inches(6.2))
footnote(s, RULE + " 절제(ablation): 공변량을 단계적으로 제거·추가해 기여를 분리하는 실험.")
notes(s, "우측 차트가 이 장의 핵심 — 정보가 쌓여도(우측으로 갈수록) BISTRO 오차는 평탄, XGB·DFM은 하락. 한은이 자체적으로 발견한 진단을 우리가 재현·원인 규명·해법(다음 장)까지 이은 흐름.")

# ============ 8. 5막: LoRA 적응 ============
s = new_slide()
header(s, "7  |  5막 · 적응학습(LoRA) 실증", "한은이 제안한 적응학습을 당사가 먼저 실증했습니다")
runs(s, 0.6, 1.45, 6.0, 4.5, [
    [("실험 (양 모델 동일 레시피, seed 3회)", True, NAVY, 12.5)],
    [("· 학습 = 과거 빈티지 경로(당시 가용 정보 + 관측 플래그,", False, INK, 10.5)],
    [("  라벨은 실제 속보치), 연 1회 fresh 재적응, release-safe", False, INK, 10.5)],
    [("· 적응 규모 = 전체 파라미터의 0.6% (LoRA rank 8)", False, INK, 10.5)],
    [("", False, INK, 5)],
    [("결과 (3-seed 예측 평균)", True, NAVY, 12.5)],
    [("· BISTRO 0.596 → 0.560 (-6.0%) — 평탄성 부분 해소", False, INK, 10.5)],
    [("· Chronos-2 0.619 → 0.564 (-8.9%)", False, INK, 10.5)],
    [("· 한은의 출력층 교체 실험(실패)과 대조 —", False, INK, 10.5)],
    [("  \"적응의 깊이\"가 관건임을 규명", True, INK, 10.5)],
    [("", False, INK, 5)],
    [("정직한 각주: BISTRO 쪽 seed 분산 큼(±0.028) — 단일 실행", False, GREY, 9.5)],
    [("불신 원칙이 여기서도 작동, 3-seed 평균만 보고", False, GREY, 9.5)],
])
s.shapes.add_picture(f"{SC}/lora_share.png", Inches(6.8), Inches(1.8), width=Inches(6.2))
footnote(s, "평가창 2021~2025(적응 데이터 확보 구간) · LoRA = 저계수 어댑터를 주의층에 삽입하는 경량 미세조정.")
notes(s, "이 장의 정치적 의미: 한은(이창훈 과장)이 8/25에 제안한 적응학습 방향을 우리가 3일 만에 첫 실증 — '제안이 옳았다'는 실증 선물로 협업 확대 논의의 주도권 확보. 실패한 변형(v2)·seed 사고도 기록으로 남겨 신뢰 유지.")

# ============ 9. 종합 성적 ============
s = new_slide()
header(s, "8  |  종합", "최종 구성 — 기준선 대비 -2.3%, 지금까지의 발견을 하나의 구성으로 정리했습니다")
s.shapes.add_picture(f"{SC}/ceo_final.png", Inches(0.5), Inches(1.6), width=Inches(7.4))
X = 8.2
sidebar(s, X, 1.7, 2.3, GREEN)
runs(s, X + 0.2, 1.7, 4.5, 2.4, [
    [("최종 운영 권고", True, NAVY, 12)],
    [("기본 = XGBoost 단독 (한은 개정판)", False, INK, 10.5)],
    [("조기 6주(지표 공백기) = FM 결합으로 교대", False, INK, 10.5)],
    [("  · 적응 이력 12분기 미만 → zero-shot 부품", False, INK, 10.5)],
    [("  · 12분기 이상 → LoRA 적응 부품 (임계 규칙)", False, INK, 10.5)],
    [("", False, INK, 4)],
    [("임계 12분기는 사전 등록값 — 초소형 적응(3~7분기)이", False, GREY, 9.5)],
    [("코로나 구간에서 역효과임을 실측으로 확인한 결과", False, GREY, 9.5)],
], sp=2.5)
sidebar(s, X, 4.25, 2.3, NAVY)
runs(s, X + 0.2, 4.25, 4.5, 2.4, [
    [("수위 (일관 유지)", True, NAVY, 12)],
    [("· -2.3%는 DM p=0.185 — 관례상 비유의,", False, INK, 10.5)],
    [("  '비열등 + 개선 방향'으로만 주장", False, INK, 10.5)],
    [("· 단, 여정 전체에서 가장 낮은 p·최다 승수(18/32)", False, INK, 10.5)],
    [("· 국면전환 v3(0.718)는 참고 기록 — 성능은 더 좋으나", False, INK, 10.5)],
    [("  수용성 문제로 보류된 연구 자산", False, INK, 10.5)],
], sp=2.5)
footnote(s, RULE)
notes(s, "차트 읽기: 왼쪽부터 성능 계보. 회색 마지막 막대(v3 0.718)가 '보류된 최고 기록'으로 남아 있음 — 대표님의 국면 직관이 숫자로는 여전히 1등이라는 점을 시각적으로 보여주는 장치이기도 함.")

# ============ 10. 협업 방식 (FDE) ============
s = new_slide()
header(s, "9  |  일하는 방식", "모든 협업을 GitHub 문서와 코드 리뷰로 진행했습니다")
runs(s, 0.6, 1.42, 12.1, 0.35, [
    [("한은이 코드와 리포트를 커밋하면, 당사가 커밋 단위로 리뷰·재검증하고 재현 가능한 스크립트와 문서로 회신하는 비동기 협업 방식이었습니다. 고객 코드를 읽는 것이 요구사항 분석이었습니다.", False, INK, 10.5)],
])
# 커밋 → 응답 대응표
runs(s, 0.6, 2.0, 5.9, 0.35, [[("이창훈 과장(한은)의 커밋", True, NAVY, 12)]])
runs(s, 6.9, 2.0, 5.8, 0.35, [[("당사의 응답", True, GREEN, 12)]])
hline(s, 0.6, 2.36, 12.13, LINE, 1.0)
pairs = [
    ("6월~  빈티지 데이터셋·평가 기준 정비", "주간 실시간 빈티지(73개 지표)·기준모형 예측·평가 그리드를 정비해 제공 — 협업의 공통 잣대를 한은이 먼저 명확히 함",
     "그 잣대 위에 하네스", "제공 그리드 그대로 재현 채점 체계 구축 — 이후 모든 실험(양측)이 같은 잣대로 비교 가능해짐"),
    ("8/4  \"코드 수정 및 산출물 저장\"", "예측단위 개정(schema v2) + 구버전 산출물 차단 검증 로직 — 기준점을 스스로 바로잡음",
     "당일 재검증", "40여 구성 전면 재채점 → 신 기준선(XGBoost 단독) 확정, 당사 종전 기록도 함께 폐기"),
    ("8/25  \"Add BISTRO nowcasting experiments...\"", "실험 코드 16종 + 테스트 + 회의자료 3.7만 줄 — 코드에서 실험 설계·한계까지 독해",
     "익일 진단 재현 → 3일 내 첫 실증", "주차별 흡수력 진단 재현·원인 규명 → 그들이 제안만 한 적응학습을 LoRA로 첫 실측"),
    ("(수신 확인도 커밋으로)", "당사 종합리포트 PDF를 한은이 자기 저장소에 커밋 — 산출물이 고객의 공식 기록에 편입",
     "산출물 = 코드", "실험 스크립트 50여 종을 결과·교훈 주석과 함께 커밋 — 실패(기각 실험)도 재현 가능한 기록으로"),
]
y = 2.48
for lt, ld, rt, rd in pairs:
    sidebar(s, 0.6, y, 0.92, NAVY)
    runs(s, 0.82, y, 5.7, 0.32, [[(lt, True, INK, 10.5)]])
    runs(s, 0.82, y + 0.32, 5.75, 0.62, [[(ld, False, GREY, 9.2)]])
    sidebar(s, 6.9, y, 0.92, GREEN)
    runs(s, 7.12, y, 2.9, 0.32, [[(rt, True, GREEN, 10.5)]])
    runs(s, 7.12, y + 0.32, 5.55, 0.62, [[(rd, False, INK, 9.2)]])
    y += 1.06
rect(s, 0.6, 6.62, 12.13, 0.42, BGGREY, None)
sidebar(s, 0.6, 6.62, 0.42, NAVY)
runs(s, 0.85, 6.68, 11.7, 0.32, [
    [("모든 과정이 커밋 기록으로 남아 언제든 재현하고 감사할 수 있습니다 — 중앙은행과 일할 때는 이 형식 자체가 신뢰를 만듭니다.", True, NAVY, 10.5)]])
footnote(s, "당사 저장소 커밋 90여 건(6~8월) · 한은 저장소는 pull 전용(푸시하지 않음) — 산출물은 문서·PDF로 전달.")
notes(s, "FDE 활동의 실체를 보여주는 장. 양방향 초상: 한은은 데이터·기준점을 명확히 하는 파트너(빈티지 정비, 자기 결함 개정, 실험 코드 공개)였고, 당사는 그 커밋을 읽고(diff에서 단위 결함·실험 설계 독해) 재현 가능한 코드로 빠르게 회신(당일 재검증, 3일 내 실증)하는 파트너. 기밀 관리 구두 포인트 — 한은 빈티지 데이터는 저장소 분리·커밋 금지 규칙 유지.")

# ============ 11. 성과와 다음 ============
s = new_slide()
header(s, "10  |  성과와 다음 단계", "검증 파트너에서 공동 연구 파트너로 — 협업의 성격이 달라졌습니다")
items = [
    ("협업 관계", "한은이 당사 자료를 자체 저장소에 수록하고, 적응학습으로의 협업 범위 확대를 먼저 제안 — 역할 분담 협의 중", GREEN),
    ("기술 자산", "실시간 검증 하네스 · 일별 빠른신호 파이프라인 · FM 절제/흡수력 진단 프레임 · LoRA 적응 러너 — 전부 재현 가능 스크립트로 관리 (50여 종)", NAVY),
    ("공동 논문", "\"사전학습 시계열 모델은 언제 가치를 더하는가\" — 기여 5건 실측 완료 + 적응학습 장(章) 추가. 한은 WP 선행 후 국제 학술지(IJF) 투고 구도", NAVY),
    ("다음 결정", "① 적응학습 역할 분담 확정(당사: Chronos-2 트랙·적응 안정화) ② 논문 집필 착수 ③ 운영 권고안의 병행 산출 개시 여부", WARM),
]
y = 1.5
for t, d, c in items:
    sidebar(s, 0.6, y, 1.1, c)
    runs(s, 0.85, y, 2.1, 0.4, [[(t, True, NAVY, 12.5)]])
    runs(s, 3.05, y, 9.6, 0.9, [[(d, False, INK, 10.5)]])
    y += 1.3
runs(s, 0.6, 6.75, 12.0, 0.3, [[("FDE로서의 3개월 — 출발점이었던 BISTRO 해석 질의에 원본 실측과 적응 개선으로 답했습니다. 코드 개정에는 당일 재검증으로, 진단에는 원인 규명으로, 제안에는 첫 실증으로 대응했습니다.", True, NAVY, 11)]])
footnote(s, RULE)
notes(s, "마지막 문장이 대표 보고의 결론 — FDE 모델의 가치 입증. 다음 결정 3건 중 대표 승인이 필요한 것은 특별히 없으며(진행 보고), 논문 저자 구성 등 회사 차원 조율이 필요해지면 별도 보고 예정.")

out = "/Users/user/vibe/bistro-lstm/docs/한은협업_여정보고_대표_2026-09-01.pptx"
prs.save(out); print("saved:", out)
