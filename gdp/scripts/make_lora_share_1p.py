# -*- coding: utf-8 -*-
"""이과장 공유 1p — LoRA 적응 실증: 귀 제안(적응학습)의 첫 실측 결과"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

F = "Apple SD Gothic Neo"
INK = RGBColor(0x26,0x26,0x26); GREY = RGBColor(0x6B,0x6B,0x6B)
LINE = RGBColor(0xC9,0xC9,0xC9); NAVY = RGBColor(0x24,0x34,0x47)
GREEN = RGBColor(0x00,0x8A,0x3E); BGGREY = RGBColor(0xF5,0xF5,0xF4)
WARM = RGBColor(0xB0,0x53,0x2F)

SC = "/private/tmp/claude-502/-Users-user-vibe-bistro-lstm/031ca4b9-003d-4cb2-afaa-06353432dc9b/scratchpad"
prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
s = prs.slides.add_slide(prs.slide_layouts[6])

def runs(x, y, w, h, lines, size=10, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp=2):
    tf = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.02); tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01); tf.margin_bottom = Inches(0.01)
    for i, rs in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp)
        for t, b, c, z in rs:
            r = p.add_run(); r.text = t
            r.font.name = F; r.font.size = Pt(z if z else size)
            r.font.bold = b; r.font.color.rgb = c if c else INK
            rp = r._r.get_or_add_rPr(); rp.append(rp.makeelement(qn('a:ea'), {'typeface': F}))
    return tf

def hline(x, y, w, color=LINE, pt=0.75):
    ln = s.shapes.add_connector(1, Inches(x), Inches(y), Inches(x+w), Inches(y))
    ln.line.color.rgb = color; ln.line.width = Pt(pt); return ln

def rect(x, y, w, h, fill=None, line=None):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False; return shp

# 헤더
runs(0.6, 0.32, 11.5, 0.3, [[("적응학습 실증 — 귀 리포트(8/25) 제안 관련  |  네이버클라우드 AX Forward Lab · 2026. 8.", True, GREY, 10.5)]])
runs(0.6, 0.6, 12.1, 0.6, [[("LoRA 적응 첫 실측 — 두 FM 모두 개선, BISTRO는 주차별 평탄성도 부분 해소", True, NAVY, 17)]])
hline(0.6, 1.22, 12.13, NAVY, 1.2)

# 설계
runs(0.6, 1.38, 12.1, 0.62, [
    [("실험 설계  ", True, NAVY, 10.5),
     ("귀 제안(과거 빈티지 전망경로를 학습표본으로 하는 제한적 추가학습)을 LoRA로 구현 — 학습단위 = 분기별 빈티지 경로(당시 이용가능 정보 + 관측여부 플래그 + 분기말 라벨은 실제 속보치), "
      "연 1회 fresh 재적응(적응 누적 없음, release-safe), 평가 2021~2025 20분기 walk-forward, seed 3개 반복.", False, INK, 10.5)],
    [("대상  ", True, NAVY, 9.5),
     ("BISTRO(공개 체크포인트, 수동 LoRA rank8 = 590K, 0.6%) · Chronos-2f(내장 LoRA fit) — 데이터·설정 완전 동일, 모델만 상이.", False, GREY, 9.5)],
])
# 좌: 차트
s.shapes.add_picture(f"{SC}/lora_share.png", Inches(0.5), Inches(2.15), width=Inches(7.55))
# 좌하: 종합표
runs(0.6, 5.5, 7.4, 0.3, [[("종합 (19주 전체 RMSE · 3-seed 예측 평균 기준)", True, NAVY, 10.5)]])
hline(0.6, 5.8, 7.4, LINE, 0.75)
runs(0.6, 5.88, 7.4, 1.1, [
    [("BISTRO   zero-shot 0.596 → ", False, INK, 9.5), ("LoRA 0.560 (-6.0%)", True, GREEN, 9.5),
     ("   ·   흡수 역행 +10.1% → +3.2%", False, INK, 9.5)],
    [("Chronos-2f   zero-shot 0.619 → ", False, INK, 9.5), ("LoRA 0.564 (-8.9%)", True, GREEN, 9.5),
     ("   ·   흡수 기울기 유지", False, INK, 9.5)],
    [("참고: XGBoost 0.518 (전체 기준선) · 2025년(사전학습과 무중복 구간)에선 BISTRO-LoRA 0.556이 XGBoost 0.567을 상회", False, GREY, 9)],
], sp=2.5)

# 우측
X, W = 8.35, 4.4
rect(X, 2.15, 0.045, 1.6, GREEN)
runs(X + 0.2, 2.15, W - 0.2, 1.65, [
    [("귀 제안 방향의 유효성이 확인됩니다", True, NAVY, 11.5)],
    [("· 출력층 교체(GDP head)로는 없던 개선이 LoRA(주의층", False, INK, 9.5)],
    [("  0.6% 적응)에서 발생 — 적응의 '깊이'가 관건이었음", False, INK, 9.5)],
    [("· 특히 BISTRO의 주차별 평탄성(귀 리포트 그림 1)이", False, INK, 9.5)],
    [("  부분 해소 — 새 빈티지 정보를 반영하기 시작", False, INK, 9.5)],
    [("· 가장 깨끗한 2025 구간에서 개선 유지 — 사전학습", False, INK, 9.5)],
    [("  중복만으로는 설명되지 않는 실질 개선", False, INK, 9.5)],
], sp=2)
rect(X, 3.9, 0.045, 1.5, WARM)
runs(X + 0.2, 3.9, W - 0.2, 1.55, [
    [("유의 — 재현성 변동이 큽니다", True, NAVY, 11.5)],
    [("· BISTRO-LoRA는 seed 간 편차가 큼 (0.549~0.602,", False, INK, 9.5)],
    [("  ±0.028) — 단일 실행 결과는 신뢰 불가", False, INK, 9.5)],
    [("· 위 수치는 3-seed 예측 평균 기준 — 운영 시에도", False, INK, 9.5)],
    [("  seed 앙상블 형태를 권장", False, INK, 9.5)],
    [("· 개선폭의 통계적 유의성은 표본(20분기) 한계로 미검정", False, INK, 9.5)],
], sp=2)
rect(X, 5.55, 0.045, 1.35, NAVY)
runs(X + 0.2, 5.55, W - 0.2, 1.4, [
    [("다음 단계 제안", True, NAVY, 11.5)],
    [("· 재현 패키지 공유 — 동일 규약 교차 검증 (귀측 BISTRO", False, INK, 9.5)],
    [("  인프라에서 독립 재현)", False, INK, 9.5)],
    [("· 공표달력·release mask 정밀화(귀측) × 적응 안정화·", False, INK, 9.5)],
    [("  seed 앙상블(당사) 분담", False, INK, 9.5)],
    [("· 공동 논문의 적응학습 장(章)으로 정리", False, INK, 9.5)],
], sp=2)

# 각주
runs(0.6, 7.12, 11.6, 0.35, [
    [("주: 실시간 빈티지 · 속보치 · 주차별 RMSE(분기 평균) · 2021~2025(20분기, 적응 데이터 확보 구간) · 낮을수록 정확. 차트는 TSFM 자체 예측 구간(w=-19~-8). "
      "BISTRO LoRA는 patch=8 고정 — zero-shot도 동일 patch 대조군 사용(공정 비교). 학습·평가에 미래 정보 미사용(release-safe).", False, GREY, 7.5)]])
runs(12.85, 7.15, 0.4, 0.25, [[("1", False, GREY, 9)]], align=PP_ALIGN.RIGHT)

out = "/Users/user/vibe/bistro-lstm/docs/LoRA적응_실증공유_1p_2026-08-27.pptx"
prs.save(out); print("saved:", out)
