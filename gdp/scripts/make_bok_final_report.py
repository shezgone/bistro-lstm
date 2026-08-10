# -*- coding: utf-8 -*-
"""한은용 종합 리포트 — 검증 총괄 · 시사점 · ICL/트랜스포머 한계 · 공동 논문 방향 (8장)"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

F = "Apple SD Gothic Neo"
INK = RGBColor(0x26,0x26,0x26); GREY = RGBColor(0x6B,0x6B,0x6B)
LINE = RGBColor(0xC9,0xC9,0xC9); NAVY = RGBColor(0x24,0x34,0x47)
GREEN = RGBColor(0x00,0x8A,0x3E); BGGREY = RGBColor(0xF3,0xF4,0xF5)
WHITE = RGBColor(0xFF,0xFF,0xFF); HL = RGBColor(0xE8,0xF3,0xEC)
WARM = RGBColor(0xB0,0x53,0x2F); BLUE = RGBColor(0x1C,0x5C,0xAB)

SC = "/private/tmp/claude-502/-Users-user-vibe-bistro-lstm/031ca4b9-003d-4cb2-afaa-06353432dc9b/scratchpad"

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
def new_slide(): return prs.slides.add_slide(prs.slide_layouts[6])

def runs(s, x, y, w, h, lines, size=10, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp=2):
    tf = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    for i, rs in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp)
        for t, b, c, z in rs:
            r = p.add_run(); r.text = t
            r.font.name = F; r.font.size = Pt(z if z else size)
            r.font.bold = b; r.font.color.rgb = c if c else INK
            rp = r._r.get_or_add_rPr(); rp.append(rp.makeelement(qn('a:ea'), {'typeface': F}))
    return tf

def hline(s, x, y, w, color=LINE, pt=0.75):
    ln = s.shapes.add_connector(1, Inches(x), Inches(y), Inches(x+w), Inches(y))
    ln.line.color.rgb = color; ln.line.width = Pt(pt); return ln

def rect(s, x, y, w, h, fill=None, line=None, round_=False):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
                             Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False; return shp

def header(s, num, title):
    runs(s, 0.55, 0.3, 12.2, 0.55, [[(num + "  ", True, GREEN, 19), (title, True, NAVY, 19)]])
    hline(s, 0.55, 0.92, 12.23, NAVY, 1.6)

def footer(s, text):
    rect(s, 0.55, 6.9, 12.23, 0.4, BGGREY, LINE)
    runs(s, 0.75, 6.96, 11.9, 0.28, [[(text, False, GREY, 8.8)]], anchor=MSO_ANCHOR.MIDDLE)

def notes(s, text):
    s.notes_slide.notes_text_frame.text = text

FOOT = "공통 검증 규약: 실시간 빈티지 · 속보치 기준 · 전망주차 w[-19,-1] 평균 RMSE · 2018Q1~2025Q4 (32개 분기) · 예측단위 개정(schema v2) 반영"

# ============ 1. 표지 ============
s = new_slide()
runs(s, 0.8, 2.4, 11.7, 0.9, [[("GDP Nowcasting 협업 종합 리포트", True, NAVY, 32)]], align=PP_ALIGN.CENTER)
runs(s, 0.8, 3.3, 11.7, 0.6, [[("검증 총괄 · 방법론 시사점 · 공동 연구 제안", False, NAVY, 19)]], align=PP_ALIGN.CENTER)
hline(s, 5.42, 4.2, 2.5, LINE, 1.0)
runs(s, 0.8, 4.4, 11.7, 0.5, [[("네이버클라우드 AX Lab  |  2026. 8.", False, GREY, 13)]], align=PP_ALIGN.CENTER)
runs(s, 0.8, 4.95, 11.7, 0.5, [[("구성: 리더보드 · 시사점 · 전망시계 진단 · 문맥학습(ICL) 방식의 한계 · 트랜스포머류의 한계와 기여 지점 · 공동 논문 방향", False, GREY, 10.5)]], align=PP_ALIGN.CENTER)
notes(s, "협업 기간(2026.6~8) 전체 검증 결과를 예측단위 개정(schema v2) 이후 기준으로 통합 정리한 자료. 모든 수치는 동일 채점 규약으로 재계산된 것이며, 개정 전 수치와의 혼용은 없음.")

# ============ 2. 리더보드 ============
s = new_slide()
header(s, "01", "검증 리더보드 — 예측단위 개정(schema v2) 이후 통합 순위")
s.shapes.add_picture(f"{SC}/bokf_leader.png", Inches(0.45), Inches(1.3), width=Inches(7.6))
X = 8.35
runs(s, X, 1.4, 4.4, 5.2, [
    [("읽는 법", True, NAVY, 11.5)],
    [("· 전 구성 동일 규약 — 직접 비교 가능", False, INK, 10)],
    [("· 개정 전 수치(예: 종전 DFM+XGBoost 0.765)는", False, INK, 10)],
    [("  단위 정합성 문제로 비교 불가 — 표에서 제외", False, INK, 10)],
    [("", False, INK, 5)],
    [("관찰 1", True, GREEN, 11)],
    [("개정 후 XGBoost 단독(0.750)이 사실상 기준선.", False, INK, 10)],
    [("DFM 결합(0.787)을 포함한 대부분의 결합이 열세", False, INK, 10)],
    [("", False, INK, 5)],
    [("관찰 2", True, GREEN, 11)],
    [("이를 수치상 상회하는 유일 구성 = 전망시계", False, INK, 10)],
    [("하이브리드(0.740) — 조기 주차만 사전학습", False, INK, 10)],
    [("모델·GBM으로 보완, 이후는 XGBoost 그대로", False, INK, 10)],
    [("", False, INK, 5)],
    [("관찰 3", True, GREEN, 11)],
    [("딥러닝 단독 최고는 Chronos-2+일별신호(0.808)", False, INK, 10)],
    [("— 기준선 미달이나 특정 조건(3장)에서 우위", False, INK, 10)],
], sp=3)
footer(s, FOOT + " · 이외 30여 개 구성(직접학습 신경망, 학습형 결합 등)은 전부 기준선 미달로 생략")
notes(s, "리더보드는 schema v2(원 단위 예측) 기준으로 전면 재계산. 전망시계 하이브리드의 -1.3%는 DM 검정 유의성 미달(p=0.41)이며, 7개 후보 구성 중 최선을 선택한 값이므로 '비열등+개선 방향' 수위로 보고. 통계적 우월 주장이 아님.")

# ============ 3. 시사점 ============
s = new_slide()
header(s, "02", "시사점 — 32개 분기 표본이 가르쳐준 세 가지")
cards = [
    ("①", "소표본에서는 단순한 고정 구조만 살아남습니다",
     ["데이터로부터 배우려는 개선책 — 결합 가중치 학습, 성과 기반 적응 가중, 잔차 보정 학습,",
      "합성데이터 미세조정 — 을 모두 동일 규약으로 검증한 결과 전부 기준선 미달.",
      "생존한 것은 사람이 고정한 단순 규칙(균등 평균, 주차별 교대)뿐 — 4개 계열 일관된 결론."]),
    ("②", "이 표본에서는 어떤 우월성도 통계적으로 입증되지 않습니다",
     ["당사 검정력 분석: 32개 분기에서는 6% 개선도 DM p=0.088. 한은 측 MCS 검정에서도",
      "비교 대상 12개 모형 전원이 5% Model Confidence Set에 잔류 — 동일한 결론.",
      "따라서 점 예측 정확도 '순위 경쟁'은 이 표본에서 판정 불가능한 게임입니다."]),
    ("③", "정확도 경쟁 대신 '조건부 가치'로 질문을 바꾸면 답이 있습니다",
     ["\"어느 모델이 이기나\"가 아니라 \"어떤 조건에서 무엇이 기여하나\"로 바꾸면 —",
      "정보 빈곤 구간(발표 19~14주 전, 경기 반등기)에서 사전학습 모델의 기여가 일관되게 관찰됩니다.",
      "이 조건부 발견이 공동 연구(6장)의 핵심 재료입니다."]),
]
y = 1.3
for num, title, body in cards:
    rect(s, 0.55, y, 12.23, 1.68, BGGREY, LINE, round_=True)
    runs(s, 0.8, y + 0.15, 0.6, 0.5, [[(num, True, GREEN, 16)]])
    runs(s, 1.45, y + 0.16, 11.0, 0.4, [[(title, True, NAVY, 12.5)]])
    runs(s, 1.45, y + 0.58, 11.1, 1.0, [[(t, False, INK, 10)] for t in body], sp=2)
    y += 1.84
footer(s, FOOT)
notes(s, "①의 '학습형 전패'는 우연이 아니라 구조적 결과 — 유효 표본 32개에서 추정하는 모든 추가 파라미터가 과적합 비용을 지불함. ②는 당사와 한은이 독립적으로 같은 결론에 도달했다는 점이 중요. ③이 본 리포트의 축.")

# ============ 4. 전망시계 진단 ============
s = new_slide()
header(s, "03", "전망시계 진단 — 사전학습 모델이 기여하는 조건")
s.shapes.add_picture(f"{SC}/bokf_horizon.png", Inches(0.45), Inches(1.35), width=Inches(6.6))
X = 7.4
runs(s, X, 1.4, 5.35, 2.5, [
    [("발견", True, NAVY, 12)],
    [("· 조기 구간(발표 19~14주 전): 해당 분기 공식지표가", False, INK, 10)],
    [("  거의 없는 정보 빈곤기 — XGBoost 우위 소멸,", False, INK, 10)],
    [("  사전학습 모델(Chronos-2f)·GBM이 근소 우위", False, INK, 10)],
    [("· 중반 이후: 지표 축적과 함께 XGBoost 독주", False, INK, 10)],
    [("· 경기 반등 분기 6개 평균: Chronos-2f 0.581 <", False, INK, 10)],
    [("  XGBoost 0.717 — 일별 금융·심리 신호의 기여", False, INK, 10)],
])
rect(s, X, 4.0, 5.35, 1.9, HL, GREEN, round_=True)
runs(s, X + 0.22, 4.15, 4.95, 1.65, [
    [("전망시계 하이브리드 (리더보드 1위 구성)", True, GREEN, 11.5)],
    [("조기 주차 = (GBM + Chronos-2f) 평균,", False, INK, 10.5)],
    [("그 외 주차 = XGBoost 그대로 — 0.740 (-1.3%)", False, INK, 10.5)],
    [("· 국면 판단 없음 — 전망 시계에 따른 조건화는", False, INK, 10)],
    [("  예측 실무의 표준 관행", False, INK, 10)],
    [("· 수위: 비열등 확인 + 개선 방향 (유의성 미달 명시)", False, INK, 10)],
], sp=3)
footer(s, FOOT + " · 반등 분기 6개: 2018Q1·2019Q2·2020Q3·2023Q1·2024Q3·2025Q2")
notes(s, "구간 경계(-14주)는 이전 fan chart 검증에서 선정의된 것을 재사용(사후 선택 아님). 다만 조기 구간에 어떤 모델 조합을 쓸지는 후보 비교로 정했으므로 선택편의 가능성을 함께 명시하는 것이 정직한 보고.")

# ============ 5. ICL 한계 ============
s = new_slide()
header(s, "04", "문맥학습(ICL) 계열의 한계 — GDP 나우캐스팅 관점")
runs(s, 0.55, 1.25, 12.2, 0.4, [[("ICL(In-Context Learning) = 모델 파라미터를 학습하지 않고, 과거 데이터를 문맥(프롬프트)으로 제시해 예측을 얻는 방식 (LLM 활용 계열 포함)", False, GREY, 10.5)]])
X1, X2 = 0.55, 6.9; W = 5.9
runs(s, X1, 1.8, W, 0.4, [[("LLM 기반 ICL 예측의 구조적 한계", True, NAVY, 12.5)]])
hline(s, X1, 2.18, W, LINE, 1.0)
runs(s, X1, 2.3, W, 4.3, [
    [("① 학습데이터 오염 (가장 결정적)", True, WARM, 11)],
    [("과거 GDP 실적·경제 뉴스가 LLM 사전학습에 이미 포함 —", False, INK, 10)],
    [("\"그 시점에 몰랐던 값\"을 모델이 알고 있을 가능성을 배제 불가.", False, INK, 10)],
    [("실시간 재현 채점 자체가 성립하지 않아 리더보드에서 제외.", True, INK, 10)],
    [("당사 예비실험도 '오염 상한'으로만 해석 (공정 점수 아님)", False, GREY, 9.3)],
    [("", False, INK, 4)],
    [("② 재현성·감사가능성", True, WARM, 11)],
    [("동일 입력에도 출력 변동, 모델 버전 변경 시 결과 비복원 —", False, INK, 10)],
    [("중앙은행 운영 요건(재현·감사·설명책임)과 정면 충돌", False, INK, 10)],
    [("", False, INK, 4)],
    [("③ 수치 정밀도", True, WARM, 11)],
    [("토큰 단위 수치 처리 특성상 소수점 정밀 회귀에 부적합", False, INK, 10)],
], sp=2)
runs(s, X2, 1.8, W, 0.4, [[("직접학습 신경망 이식판의 실측 (참고)", True, NAVY, 12.5)]])
hline(s, X2, 2.18, W, LINE, 1.0)
runs(s, X2, 2.3, W, 4.3, [
    [("사내 어텐션 LSTM 계열(BISTRO 유사 구조)을 당사가 GDP", False, INK, 10)],
    [("과제용으로 재구현해 동일 규약으로 채점한 결과:", False, INK, 10)],
    [("", False, INK, 3)],
    [("· 단독 1.268 — 기준선(0.750)의 약 1.7배 오차", False, INK, 10.5)],
    [("· 경기 반등 분기 2.401 — 학습분포 밖 급변기 일반화 실패", False, INK, 10.5)],
    [("· XGBoost와 오차 상관 0.847 — 동일 패널 소비로", False, INK, 10.5)],
    [("  결합 다양성 부재 (모든 결합 배합에서 단조 열화)", False, INK, 10.5)],
    [("· 원인: 월 140행 표본으로 다수 파라미터 직접학습", False, INK, 10.5)],
    [("", False, INK, 5)],
    [("※ 유의: 본 결과는 당사 재구현 이식판에 대한 동일 규약", True, GREY, 9.5)],
    [("   판정이며, 원 연구·타 과제(물가 등)에서의 성능에 대한", False, GREY, 9.5)],
    [("   판정이 아님. 원본도 동일 규약 채점 시 공정 비교 가능.", False, GREY, 9.5)],
], sp=2)
footer(s, FOOT)
notes(s, "좌측: LLM ICL은 성능 이전에 평가 타당성(오염) 문제로 실시간 나우캐스팅 채점이 성립하지 않음 — 이것이 리더보드 제외 사유. 우측: 이식판 실측은 아키텍처 계열의 소표본 한계에 대한 증거이며, 특정 연구에 대한 평가가 아님을 명시적으로 구분.")

# ============ 6. 트랜스포머류 한계 ============
s = new_slide()
header(s, "05", "트랜스포머류(파운데이션 모델)가 점 예측에 기여하기 어려운 이유")
reasons = [
    ("소표본 정형 데이터는 GBDT의 영역", "월 140행·지표 34종 조건에서 튜닝된 트리 모델 우위는 문헌 일반 결과(Grinsztajn et al., NeurIPS 2022)와 당사 실측(직접학습 신경망 전멸)이 일치"),
    ("사전학습 지식 ≠ 한국 거시 구조", "파운데이션 모델이 배운 것은 시계열의 보편 문법(추세·계절성) — 지표 간 인과·한국 특수성은 없음. 공변량 소화력도 XGBoost(34종 비선형)에 미달"),
    ("정보 상한", "합성데이터로 추가 미세조정해도 예측이 거의 불변(상관 0.990) — 병목은 모델·표본량이 아니라 입력 패널이 가진 정보 자체"),
    ("분포 출력의 과소산포", "원시 분위수의 실측 커버리지 49%(명목 80%) — 그대로 사용 불가, 사후 보정 필수"),
]
y = 1.3
for i, (t, d) in enumerate(reasons):
    rect(s, 0.55, y, 7.5, 1.22, BGGREY, LINE, round_=True)
    runs(s, 0.78, y + 0.12, 0.5, 0.4, [[(str(i+1), True, WARM, 14)]])
    runs(s, 1.35, y + 0.13, 6.5, 0.35, [[(t, True, NAVY, 11.5)]])
    runs(s, 1.35, y + 0.5, 6.55, 0.65, [[(d, False, INK, 9.5)]])
    y += 1.36
rect(s, 8.3, 1.3, 4.48, 5.4, HL, GREEN, round_=True)
runs(s, 8.55, 1.5, 4.0, 5.0, [
    [("그럼에도 기여가 확인된 지점", True, GREEN, 12.5)],
    [("", False, INK, 4)],
    [("① 정보 빈곤 구간", True, INK, 11)],
    [("발표 19~14주 전 — 공식지표 공백기에", False, INK, 10)],
    [("세계 지식 + 일별 신호로 보완 (3장)", False, INK, 10)],
    [("", False, INK, 4)],
    [("② 경기 반등기", True, INK, 11)],
    [("반등 6개 분기 평균 0.581 — 전 구성 중", False, INK, 10)],
    [("최강. 주가·환율·심리 원지수가 공식", False, INK, 10)],
    [("지표의 발표 시차를 메움", False, INK, 10)],
    [("", False, INK, 4)],
    [("③ 불확실성 정량화 재료", True, INK, 11)],
    [("보정(conformal) 결합 시 예측구간의", False, INK, 10)],
    [("품질 개선 확인 — 연구 재료로 유효", False, INK, 10)],
    [("", False, INK, 4)],
    [("요약: 대체재가 아니라 조건부 보완재", True, GREEN, 11)],
], sp=2.5)
footer(s, FOOT)
notes(s, "핵심 메시지: '트랜스포머가 안 된다'가 아니라 '튜닝된 GBDT가 지배하는 조건과, 사전학습 지식이 기여하는 조건이 다르다'. 좌측 4개는 실측 근거가 각각 존재. 우측 조건부 기여가 공동 논문의 핵심 주장으로 연결.")

# ============ 7. 공동 논문 ============
s = new_slide()
header(s, "06", "공동 논문 제안 — \"사전학습 시계열 모델은 언제 가치를 더하는가\"")
runs(s, 0.55, 1.25, 12.2, 0.55, [
    [("프레임: 우월성 주장이 아닌 조건부 가치 규명 — ", False, INK, 11),
     ("\"When do time-series foundation models add value? A real-time evaluation of GDP nowcasting\"", True, NAVY, 11)],
])
X1, W1 = 0.55, 6.4
runs(s, X1, 1.9, W1, 0.4, [[("기여 (모두 실측 완료)", True, NAVY, 12)]])
hline(s, X1, 2.26, W1, LINE, 1.0)
runs(s, X1, 2.38, W1, 3.6, [
    [("1. ", True, GREEN, 10.5), ("2세대 TSFM의 실시간 빈티지 평가 최초 사례", True, INK, 10.5)],
    [("   주 단위 19회 예측 그리드 — 기존 문헌은 대부분 사후 확정 데이터", False, GREY, 9.3)],
    [("2. ", True, GREEN, 10.5), ("전망시계 조건부 보완성", True, INK, 10.5)],
    [("   정보 빈곤 구간·반등기에서만 기여 — 메커니즘 있는 조건부 발견", False, GREY, 9.3)],
    [("3. ", True, GREEN, 10.5), ("무개정 일별 신호의 식별 깨끗한 효과 (0.836→0.808)", True, INK, 10.5)],
    [("4. ", True, GREEN, 10.5), ("TSFM 분위수 과소산포 실측 + 사후보정 처방", True, INK, 10.5)],
    [("5. ", True, GREEN, 10.5), ("소표본 결합 퍼즐 — 학습형 전패·고정 구조 생존 + 검정력 분석", True, INK, 10.5)],
    [("   MCS 전원 생존 결과와 함께 '평가 축 재설계' 논의로 승격", False, GREY, 9.3)],
])
rect(s, X1, 5.6, W1, 1.1, BGGREY, LINE, round_=True)
runs(s, X1 + 0.2, 5.72, W1 - 0.4, 0.9, [
    [("역할 분담(안)  ", True, NAVY, 10.5),
     ("한은: 실시간 빈티지·도메인 검증·정책 함의  |  당사: 모델·평가 인프라·실험 재현 패키지", False, INK, 10)],
    [("발표 경로(안)  ", True, NAVY, 10.5),
     ("한은 Working Paper 선행 공개 → International Journal of Forecasting 투고", False, INK, 10)],
], sp=3)
X2, W2 = 7.35, 5.43
runs(s, X2, 1.9, W2, 0.4, [[("표본 한계의 정면 돌파 (Future Work 아님 — 설계에 포함)", True, NAVY, 11.5)]])
hline(s, X2, 2.26, W2, LINE, 1.0)
runs(s, X2, 2.38, W2, 3.0, [
    [("· 다국가 복제", True, INK, 10.5)],
    [("  미국 ALFRED(세인트루이스 연준 실시간 DB, 공개) ·", False, INK, 10)],
    [("  유로존 ECB 실시간 DB로 동일 프레임 재현 —", False, INK, 10)],
    [("  3개 지역 일관 결과로 소표본 반론에 선제 대응", False, INK, 10)],
    [("", False, INK, 4)],
    [("· 사전등록 라이브 추적", True, INK, 10.5)],
    [("  하이브리드 구성을 고정·공개 후 향후 분기를 실전", False, INK, 10)],
    [("  축적 — 선택편의 없는 표본외 증거 확보", False, INK, 10)],
    [("", False, INK, 4)],
    [("· 유의성 논의는 본문에서 정면 처리", True, INK, 10.5)],
    [("  검정력 분석 + MCS를 '평가 방법론 기여'로 배치 —", False, INK, 10)],
    [("  약점 고백이 아니라 논문 주장의 일부", False, INK, 10)],
], sp=2)
footer(s, "데이터 공개 방식(코드 공개 + 빈티지 데이터 접근 절차)은 초기 합의 필요 · 기준선 개정 경위의 서술 수위는 공동저자 협의 사항")
notes(s, "제목·프레임의 핵심은 '이긴다'가 아니라 '언제 기여하는가'. 이 구도에서는 유의성 미달이 약점이 아니라 5번 기여(평가 방법론)의 일부가 됨. 다국가 복제는 리뷰어 1순위 예상 공격(단일국 32분기)에 대한 선제 설계.")

# ============ 8. 요약 ============
s = new_slide()
header(s, "07", "요약 및 다음 단계")
rows = [
    ("점 예측", "XGBoost 단독(개정 후)이 기준선 — 인정. 전망시계 하이브리드(0.740)는 비열등+개선 방향의 참고 옵션", "필요 시 병행 산출"),
    ("문맥학습(ICL)", "LLM 계열은 오염 문제로 실시간 채점 자체가 비성립 · 직접학습 이식판은 소표본 한계 실측", "채택 비권고"),
    ("트랜스포머류", "점 예측 대체재 아님 — 정보 빈곤 구간·반등기의 조건부 보완재 + 불확실성 재료", "논문 핵심 재료"),
    ("공동 논문", "\"언제 가치를 더하는가\" 프레임 · 기여 5건 실측 완료 · 다국가 복제 설계 포함", "구도 협의 요청"),
]
y = 1.35
for a, bb, c in rows:
    rect(s, 0.55, y, 12.23, 1.12, BGGREY, LINE, round_=True)
    runs(s, 0.8, y + 0.13, 1.9, 0.4, [[(a, True, NAVY, 12)]])
    runs(s, 2.8, y + 0.13, 7.6, 0.85, [[(bb, False, INK, 10.5)]])
    runs(s, 10.6, y + 0.35, 2.0, 0.4, [[(c, True, GREEN, 10.5)]], align=PP_ALIGN.CENTER)
    y += 1.28
runs(s, 0.55, 6.5, 12.2, 0.35, [[("모든 실험은 재현 가능한 스크립트로 관리 중이며, 논문 착수 결정 시 재현 패키지 형태로 정리해 공유 가능합니다.", False, GREY, 10)]])
footer(s, FOOT)
notes(s, "다음 회의 의사결정 요청 사항: ① 공동 논문 구도(제목 프레임·역할 분담·발표 경로) 협의 착수 여부 ② 하이브리드 병행 산출 여부(선택). 나머지는 보고 사항.")

out = "/Users/user/vibe/bistro-lstm/docs/GDP_Nowcasting_종합리포트_한은_2026-08-10.pptx"
prs.save(out); print("saved:", out)
