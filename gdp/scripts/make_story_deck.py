# -*- coding: utf-8 -*-
"""상사 보고용 스토리 덱 — 세 명의 예측가 이야기 (5장)"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

F = "Apple SD Gothic Neo"
INK = RGBColor(0x26,0x26,0x26); GREY = RGBColor(0x6B,0x6B,0x6B)
LINE = RGBColor(0xC9,0xC9,0xC9); NAVY = RGBColor(0x24,0x34,0x47)
GREEN = RGBColor(0x00,0x8A,0x3E); BGGREY = RGBColor(0xF3,0xF4,0xF5)
WHITE = RGBColor(0xFF,0xFF,0xFF); HL = RGBColor(0xE8,0xF3,0xEC)
WARM = RGBColor(0xB0,0x53,0x2F); BLUE = RGBColor(0x1C,0x5C,0xAB)

SCRATCH = "/private/tmp/claude-502/-Users-user-vibe-bistro-lstm/031ca4b9-003d-4cb2-afaa-06353432dc9b/scratchpad"

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

def new_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def runs(s, x, y, w, h, lines, size=10, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp=2):
    tf = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    for i, rs in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp)
        for t, b, c, z in rs:
            r = p.add_run(); r.text = t
            r.font.name = F; r.font.size = Pt(z if z else size)
            r.font.bold = b; r.font.color.rgb = c if c else INK
            rp = r._r.get_or_add_rPr(); rp.append(rp.makeelement(qn('a:ea'), {'typeface': F}))
    return tf

def hline(s, x, y, w, color=LINE, pt=0.75):
    ln = s.shapes.add_connector(1, Inches(x), Inches(y), Inches(x+w), Inches(y))
    ln.line.color.rgb = color; ln.line.width = Pt(pt); return ln

def rect(s, x, y, w, h, fill=None, line=None, round_=False):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
                             Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False; return shp

def header(s, num, title, sub=None):
    runs(s, 0.55, 0.3, 12.2, 0.55, [[(num + "  ", True, GREEN, 20), (title, True, NAVY, 20)]])
    hline(s, 0.55, 0.92, 12.23, NAVY, 1.6)
    if sub:
        runs(s, 0.55, 1.0, 12.2, 0.4, [[(sub, False, GREY, 11.5)]])

def footer(s, text):
    rect(s, 0.55, 6.85, 12.23, 0.42, BGGREY, LINE)
    runs(s, 0.75, 6.91, 11.9, 0.3, [[(text, False, GREY, 9)]], anchor=MSO_ANCHOR.MIDDLE)

# ================= 슬라이드 1: 표지 =================
s = new_slide()
rect(s, 0, 0, 13.333, 7.5, WHITE)
runs(s, 0.8, 2.05, 11.7, 0.9, [[("🏛️   🔧   🛰️", False, INK, 40)]], align=PP_ALIGN.CENTER)
runs(s, 0.8, 3.0, 11.7, 0.9, [[("세 명의 예측가 이야기", True, NAVY, 40)]], align=PP_ALIGN.CENTER)
runs(s, 0.8, 3.95, 11.7, 0.5, [[("— GDP 나우캐스트가 국면 판단 없이 신기록을 낸 비결 —", False, GREY, 16)]], align=PP_ALIGN.CENTER)
runs(s, 0.8, 4.75, 11.7, 0.5, [[("원로 이코노미스트, 꼼꼼한 조수, 그리고 시세판을 보는 컨설턴트", False, INK, 13)]], align=PP_ALIGN.CENTER)
hline(s, 5.17, 5.5, 3.0, LINE, 1.0)
runs(s, 0.8, 5.7, 11.7, 0.4, [[("AX Lab · 2026. 07  |  한국은행 GDP Nowcasting 협업", False, GREY, 11)]], align=PP_ALIGN.CENTER)

# ================= 슬라이드 2: 등장인물 =================
s = new_slide()
header(s, "01", "등장인물 — 같은 문제를 푸는 세 개의 다른 눈")
cards = [
    ("🏛️", "원로 이코노미스트", "DFM (현행 시스템의 중심)", [
        ("정통 요인모형. 공식지표 34종을 종합해", False),
        ("월별 GDP 흐름의 '초안'을 작성", False),
        ("보는 것: 물가·수출·생산·심리 등 34종", True),
        ("못 보는 것: 주가·환율 시세판 (규정 외)", True),
        ("약점: 급변기에 반 박자 느림", False),
    ]),
    ("🔧", "꼼꼼한 조수", "XGBoost (머신러닝 보정)", [
        ("원로의 초안을 받아 과거 12년 패턴으로", False),
        ("비선형 보정 — 원로와 콤비로 현행 최고", False),
        ("보는 것: 원로와 같은 34종 + 원로의 초안", True),
        ("못 보는 것: 역시 시세판은 못 봄", True),
        ("약점: 가끔 혼자 지나치게 비관·낙관", False),
    ]),
    ("🛰️", "외국인 컨설턴트", "Chronos-2f (사전학습 AI, 딥러닝)", [
        ("한국 데이터로 공부한 적 없음 — 전 세계", False),
        ("수백만 개 시계열 패턴을 미리 학습(zero-shot)", False),
        ("보는 것: 원로의 초안 + 유일하게 시세판", True),
        ("  (일별 KOSPI·원달러 + 심리 원지수)", None),
        ("약점: 평온기엔 원로+조수 콤비보다 부정확", False),
    ]),
]
for i, (emoji, name, role, bullets) in enumerate(cards):
    x = 0.55 + i * 4.18
    rect(s, x, 1.35, 3.95, 3.95, BGGREY, LINE, round_=True)
    runs(s, x + 0.25, 1.6, 3.5, 0.7, [[(emoji, False, INK, 30)]])
    runs(s, x + 0.25, 2.35, 3.5, 0.4, [[(name, True, NAVY, 14)]])
    runs(s, x + 0.25, 2.72, 3.5, 0.35, [[(role, False, GREY, 10.5)]])
    hline(s, x + 0.25, 3.12, 3.45, LINE, 0.75)
    lines = []
    for t, b in bullets:
        if b is None:
            lines.append([(t, False, GREY, 9.5)])
        else:
            lines.append([("· ", False, GREY, 9.5), (t, b, INK if not b else NAVY, 9.5)])
    runs(s, x + 0.25, 3.25, 3.5, 2.4, lines, sp=4)
runs(s, 0.55, 5.6, 12.2, 0.6, [[
    ("성적표 (8년 32분기 실시간 재현 오차, 낮을수록 정확)   ", True, NAVY, 11),
    ("원로 단독 0.865   ·   원로+조수 0.765 ", False, INK, 11),
    ("← 현행 최고", True, GREEN, 11),
    ("   ·   컨설턴트 단독 0.808", False, INK, 11)]])
footer(s, "모두 동일 조건: 한은 실시간 빈티지, 속보치 기준, 전망주차 w[-19,-1] 평균 RMSE, 2018Q1~2025Q4")

# ================= 슬라이드 3: 사건 =================
s = new_slide()
header(s, "02", "사건 — 2020년 여름, 시세판은 알고 있었다")
runs(s, 0.55, 1.3, 5.5, 4.8, [
    [("코로나 직후, 3분기 반등을 맞힐 수 있는가", True, NAVY, 13)],
    [("", False, INK, 6)],
    [("· 공식 심리지표(ESI)는 4월 55.7로 바닥을 찍고", False, INK, 11)],
    [("  여름까지도 채 회복하지 못한 상태", False, INK, 11)],
    [("· 그런데 주식시장은 4월부터 넉 달 연속 강한 반등", False, INK, 11)],
    [("  — 주가·환율은 발표를 기다릴 필요 없이 당일 확정", False, INK, 11)],
    [("", False, INK, 6)],
    [("발표 16주 전, 각자의 예측은", True, NAVY, 13)],
    [("", False, INK, 6)],
    [("· 원로(공식지표만): ", False, INK, 11), ("-0.5%", True, WARM, 12), ("  \"아직 침체입니다\"", False, GREY, 10.5)],
    [("· 컨설턴트(시세판 포함): ", False, INK, 11), ("-0.27%", True, BLUE, 12), ("  \"반등 조짐이 보입니다\"", False, GREY, 10.5)],
    [("· 실제 속보치: ", False, INK, 11), ("+1.9%", True, GREEN, 13)],
    [("", False, INK, 6)],
    [("모두 틀렸지만 — 시세판을 본 쪽이 덜 틀렸고,", False, INK, 11.5)],
    [("발표가 가까워질수록 그 격차는 벌어졌습니다.", False, INK, 11.5)],
])
s.shapes.add_picture(f"{SCRATCH}/story_c1.png", Inches(6.3), Inches(1.35), width=Inches(6.5))
footer(s, "반등 분기 6개(2018Q1·2019Q2·2020Q3·2023Q1·2024Q3·2025Q2) 평균에서도 컨설턴트(0.581)가 원로(0.621)보다 정확 — 시세판 정보의 가치")

# ================= 슬라이드 4: 반전 =================
s = new_slide()
header(s, "03", "반전 — 원로와 컨설턴트, 둘만 섞었더니 오히려 나빠졌다")
s.shapes.add_picture(f"{SCRATCH}/story_c2.png", Inches(0.45), Inches(1.4), width=Inches(6.2))
X = 7.0
runs(s, X, 1.35, 5.75, 1.0, [
    [("\"좋은 둘을 평균하면 더 좋아지겠지?\" — 아니었습니다", True, NAVY, 13)],
    [("두 사람 평균의 8년 성적: 0.808 → ", False, INK, 11.5), ("0.828 (후퇴)", True, WARM, 12)],
])
rect(s, X, 2.35, 5.75, 1.4, BGGREY, LINE, round_=True)
runs(s, X + 0.25, 2.5, 5.3, 1.2, [
    [("이유: 컨설턴트는 백지에서 시작하지 않는다", True, NAVY, 11.5)],
    [("컨설턴트가 받아보는 '과거 GDP 흐름' 자체가 원로가 만든", False, INK, 10.5)],
    [("월별 초안입니다. 즉 컨설턴트의 의견에는 이미 원로가", False, INK, 10.5)],
    [("절반쯤 들어있습니다.", True, INK, 10.5)],
], sp=3)
runs(s, X, 4.05, 5.75, 2.3, [
    [("그 상태에서 원로와 또 평균하면:", True, INK, 11)],
    [("", False, INK, 4)],
    [("· 사실상 ", False, INK, 10.5), ("원로 의견 1.5표 + 시세판 정보 0.5표", True, NAVY, 10.5)],
    [("· 유일하게 새로웠던 시세판 정보가 희석되고,", False, INK, 10.5)],
    [("· 원로가 틀린 구간에서는 그 오차까지 전염됩니다", False, INK, 10.5)],
    [("", False, INK, 4)],
    [("→ 회의에서 같은 사람 말을 두 번 듣는 것과 같습니다.", True, WARM, 11)],
])
footer(s, "원로 단독+컨설턴트 50:50 평균 = 0.828 (컨설턴트 단독 0.808보다 열세) — 정보 중복과 오차 전염이 원인")

# ================= 슬라이드 5: 결말 =================
s = new_slide()
header(s, "04", "결말 — 세 명이 각자 한 표씩, '균등 위원회'")
s.shapes.add_picture(f"{SCRATCH}/story_c3.png", Inches(0.4), Inches(1.5), width=Inches(6.9))
X = 7.55
runs(s, X, 1.35, 5.2, 2.6, [
    [("(원로 + 조수 + 컨설턴트) ÷ 3 = ", True, NAVY, 12.5), ("0.746", True, GREEN, 14)],
    [("", False, INK, 4)],
    [("· 현행 최고(0.765) 대비 ", False, INK, 11), ("-2.5%", True, GREEN, 11.5), (" — 국면 판단 없이 달성", False, INK, 11)],
    [("· 32개 분기 중 ", False, INK, 11), ("20개 분기에서 개선", True, INK, 11), (" (코로나 시기 제외에도 유효)", False, GREY, 10)],
    [("· 반등 분기 오차 0.815 → ", False, INK, 11), ("0.715", True, INK, 11), (" — 시세판 정보가 제몫", False, GREY, 10)],
    [("· 가중치 튜닝 없음 — 단순 균등 1표씩이 최적", False, INK, 11)],
], sp=4)
rect(s, X, 4.05, 5.2, 1.35, HL, GREEN, round_=True)
runs(s, X + 0.22, 4.2, 4.8, 1.15, [
    [("왜 위원회는 이기는가", True, GREEN, 11.5)],
    [("세 사람은 서로 다른 순간에, 서로 다른 방향으로 틀립니다.", False, INK, 10.5)],
    [("조수가 혼자 비관할 때 컨설턴트가 되돌리고, 컨설턴트가", False, INK, 10.5)],
    [("평온기에 흔들릴 때 원로+조수가 잡아줍니다.", False, INK, 10.5)],
], sp=3)
runs(s, X, 5.6, 5.2, 0.9, [
    [("유의사항: 개선폭 2.5%는 32분기 표본으로는 통계적 유의성 미달", False, GREY, 9.5)],
    [("(p=0.37) — 보고 수위는 \"기존 대비 비열등 + 개선 방향\"이 정직합니다.", False, GREY, 9.5)],
], sp=2)
footer(s, "다음 갈림길: 이 결과는 '중심선 교체' 옵션 — 예측 불변이 전제인 fan chart 제안과 별개 트랙으로, 한은 수용성에 따라 선택 제시 가능")

out = "/Users/user/vibe/bistro-lstm/docs/GDP_앙상블_스토리_세예측가_2026-07-29.pptx"
prs.save(out); print("saved:", out)
