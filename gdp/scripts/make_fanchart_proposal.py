# -*- coding: utf-8 -*-
"""한은 제안 1장 — GDP 나우캐스트 Fan Chart (예측 불변 레이어)"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

F = "Apple SD Gothic Neo"
INK = RGBColor(0x26,0x26,0x26); GREY = RGBColor(0x6B,0x6B,0x6B)
LINE = RGBColor(0xC9,0xC9,0xC9); NAVY = RGBColor(0x24,0x34,0x47)
GREEN = RGBColor(0x00,0x8A,0x3E); BGGREY = RGBColor(0xF3,0xF4,0xF5)
WHITE = RGBColor(0xFF,0xFF,0xFF); HL = RGBColor(0xE8,0xF3,0xEC)

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
s = prs.slides.add_slide(prs.slide_layouts[6])

def runs(tf, lines, size=10, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp=2):
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    for i, rs in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp)
        for t, b, c, z in rs:
            r = p.add_run(); r.text = t
            r.font.name = F; r.font.size = Pt(z if z else size)
            r.font.bold = b; r.font.color.rgb = c if c else INK
            rp = r._r.get_or_add_rPr(); rp.append(rp.makeelement(qn('a:ea'), {'typeface': F}))

def tb(x, y, w, h): return s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
def hline(x, y, w, color=LINE, pt=0.75):
    ln = s.shapes.add_connector(1, Inches(x), Inches(y), Inches(x+w), Inches(y))
    ln.line.color.rgb = color; ln.line.width = Pt(pt); return ln
def rect(x, y, w, h, fill=None, line=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb = line; sh.line.width = Pt(0.75)
    sh.shadow.inherit = False; return sh

# 헤더
runs(tb(0.55, 0.26, 11.0, 0.55), [[("GDP 나우캐스트 Fan Chart — 예측은 그대로, 불확실성을 입힙니다", True, NAVY, 21)]])
runs(tb(11.3, 0.36, 1.5, 0.4), [[("AX Lab | 2026. 07.", False, GREY, 10)]], align=PP_ALIGN.RIGHT)
hline(0.55, 0.88, 12.23, NAVY, 1.6)
runs(tb(0.55, 0.96, 12.2, 0.42), [[
    ("현행 점 전망(DFM+XGBoost)은 한 글자도 바꾸지 않고, 주차별 예측구간(50·80%)을 부여 — ", False, INK, 11.5),
    ("실시간 검증 커버리지 81~87% (명목 80%)", True, GREEN, 12)]])

COLY = 1.52
def col_header(x, w, num, title):
    runs(tb(x, COLY, w, 0.34), [[(num+"  ", True, GREEN, 12.5), (title, True, NAVY, 12.5)]])
    hline(x, COLY+0.36, w, LINE, 1.0)

# 01 방법 (좌)
x, w = 0.55, 3.7
col_header(x, w, "01", "방법 — 예측 불변 부가 레이어")
runs(tb(x, COLY+0.46, w, 4.5), [
    [("중심선 = 현행 점 전망 그대로", True, INK, 10)],
    [("· 모형·스킴 변경 없음. 구간만 얹음", False, INK, 9.3)],
    [("구간 = 과거 실시간 오차의 경험분포", True, INK, 10)],
    [("· 컨포멀 예측: 분포 가정 없이 통계적", False, INK, 9.3)],
    [("  보장, 전망주차 구간별로 폭 차등", False, INK, 9.3)],
    [("· 학습·보정에 미래정보 미사용", False, INK, 9.3)],
    [("  (확장창, 발표 전 분기 제외)", False, GREY, 8.8)],
    [("조기 주차 정밀화 (딥러닝 기여)", True, INK, 10)],
    [("· 발표 19~14주 전 구간은 사전학습", False, INK, 9.3)],
    [("  시계열 모델(Chronos-2)의 분포를", False, INK, 9.3)],
    [("  같은 방식으로 보정해 사용", False, INK, 9.3)],
    [("· 전망 시계별 조건화 — 국면 판단이", False, INK, 9.3)],
    [("  아닌 표준 관행", False, GREY, 8.8)],
])

runs(tb(0.55, COLY+3.45, 3.7, 1.4), [
    [("도입 형태", True, NAVY, 10)],
    [("· 주간 전망 보고서에 그림 1장 추가가 전부", False, INK, 9.3)],
    [("· 기존 파이프라인 무수정 (산출물만 소비)", False, INK, 9.3)],
    [("· 전 코드 재현 가능 형태로 이관", False, INK, 9.3)],
])

# 02 검증 (중)
x, w = 4.55, 3.6
col_header(x, w, "02", "실시간 검증 (2020~2025, 23분기)")
data = [("지표", "구간 성능", ""),
        ("커버리지(명목 80%)", "81%", "전체"),
        ("〃 (COVID 3분기 제외)", "87%", ""),
        ("평균 구간 폭", "1.86%p", "폭도 축소"),
        ("구간 점수(Winkler)", "−5.5%", "DL 결합 효과")]
t = s.shapes.add_table(5, 3, Inches(x), Inches(COLY+0.5), Inches(w), Inches(1.6)).table
t.columns[0].width = Inches(1.85); t.columns[1].width = Inches(0.85); t.columns[2].width = Inches(0.9)
for r in range(5):
    for c in range(3):
        cell = t.cell(r, c)
        cell.margin_left = Inches(0.04); cell.margin_right = Inches(0.03)
        cell.margin_top = Inches(0.012); cell.margin_bottom = Inches(0.012)
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY if r==0 else (HL if r in (1,2) else WHITE)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if c==0 else PP_ALIGN.CENTER
        run = p.add_run(); run.text = data[r][c]
        run.font.name = F; run.font.size = Pt(8.8); run.font.bold = (r==0 or (c==1 and r in (1,2)))
        run.font.color.rgb = WHITE if r==0 else INK
        rp = run._r.get_or_add_rPr(); rp.append(rp.makeelement(qn('a:ea'), {'typeface': F}))
runs(tb(x, COLY+2.25, w, 2.6), [
    [("의의", True, INK, 10)],
    [("· 현행 시스템은 점추정만 제공 — 전망의", False, INK, 9.3)],
    [("  불확실성이 보고서에 수치로 없음", False, INK, 9.3)],
    [("· 주요 중앙은행 관행(BoE Inflation Report", False, INK, 9.3)],
    [("  fan chart)과 동일한 커뮤니케이션 형식", False, INK, 9.3)],
    [("· 발표가 다가올수록 구간이 수축 —", False, INK, 9.3)],
    [("  \"확신의 정도\"가 주차별로 표현됨", False, INK, 9.3)],
    [("· 다음 단계: 한은 뉴스심리지수(NSI) 결합", True, INK, 9.3)],
    [("  으로 조기 구간 정밀화, 전망 변동", False, INK, 9.3)],
    [("  자동 설명 리포트와 통합", False, INK, 9.3)],
])

# 03 그림 (우)
x, w = 8.45, 4.35
col_header(x, w, "03", "예시 — 2025년 2·3분기 (실시간 재현)")
s.shapes.add_picture("/Users/user/vibe/gdp-nowcasting-renewal/output/fanchart_demo.png",
                     Inches(x-0.25), Inches(COLY+0.7), width=Inches(w+0.55))
runs(tb(x, COLY+2.75, w, 0.7), [
    [("점선 = 이후 발표된 실제 속보치. 두 분기 모두 80% 구간 내에서", False, GREY, 8.8)],
    [("실현 — 반등 국면(2025Q2)에서도 구간이 리스크를 사전에 포괄", False, GREY, 8.8)],
])

# 하단
rect(0.55, 6.7, 12.23, 0.5, BGGREY, LINE)
runs(tb(0.75, 6.78, 11.9, 0.36), [[
    ("검증 기준   ", True, NAVY, 9.5),
    ("한은 제공 실시간 빈티지, 속보치 기준, 전망주차 w[−19,−1] · 구간 산출에 쓰인 오차는 해당 분기 이전 발표분만 사용 (release-safe)", False, INK, 9.3)]], anchor=MSO_ANCHOR.MIDDLE)

out = "/Users/user/vibe/bistro-lstm/docs/GDP_FanChart_제안_1p_2026-07-27.pptx"
prs.save(out); print("saved:", out)
