# -*- coding: utf-8 -*-
"""한은용 종합 리포트 v2 — 실무자 스타일 디자인 (메시지형 제목·괘선·좌측 색 바·각주·쪽번호)"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

F = "Apple SD Gothic Neo"
INK = RGBColor(0x26,0x26,0x26); GREY = RGBColor(0x6B,0x6B,0x6B)
LINE = RGBColor(0xC9,0xC9,0xC9); NAVY = RGBColor(0x24,0x34,0x47)
GREEN = RGBColor(0x00,0x8A,0x3E); BGGREY = RGBColor(0xF5,0xF5,0xF4)
WHITE = RGBColor(0xFF,0xFF,0xFF); HL = RGBColor(0xEE,0xF5,0xF0)
WARM = RGBColor(0xB0,0x53,0x2F); BLUE = RGBColor(0x1C,0x5C,0xAB)

SC = "/private/tmp/claude-502/-Users-user-vibe-bistro-lstm/031ca4b9-003d-4cb2-afaa-06353432dc9b/scratchpad"

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
PAGE = [0]
def new_slide():
    PAGE[0] += 1
    return prs.slides.add_slide(prs.slide_layouts[6])

def runs(s, x, y, w, h, lines, size=10, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp=2):
    tf = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.02); tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01); tf.margin_bottom = Inches(0.01)
    for i, rs in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp)
        for t, b, c, z in rs:
            r = p.add_run(); r.text = t
            r.font.name = F; r.font.size = Pt(z if z else size)
            r.font.bold = b; r.font.color.rgb = c if c else INK
            rp = r._r.get_or_add_rPr(); rp.append(rp.makeelement(qn('a:ea'), {'typeface': F}))
    return tf

def hline(s, x, y, w, color=LINE, pt=0.75):
    ln = s.shapes.add_connector(1, Inches(x), Inches(y), Inches(x+w), Inches(y))
    ln.line.color.rgb = color; ln.line.width = Pt(pt); return ln

def rect(s, x, y, w, h, fill=None, line=None):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False; return shp

def header(s, tag, msg):
    """실무 스타일: 작은 섹션 태그 + 메시지형 제목 + 가는 밑줄"""
    runs(s, 0.6, 0.32, 11.5, 0.3, [[(tag, True, GREY, 10.5)]])
    runs(s, 0.6, 0.6, 12.1, 0.6, [[(msg, True, NAVY, 17)]])
    hline(s, 0.6, 1.22, 12.13, NAVY, 1.2)

def footnote(s, text):
    runs(s, 0.6, 7.06, 10.8, 0.3, [[("주: " + text, False, GREY, 8)]])
    runs(s, 12.35, 7.06, 0.5, 0.3, [[(str(PAGE[0]), False, GREY, 9)]], align=PP_ALIGN.RIGHT)

def sidebar(s, x, y, h, color=NAVY):
    rect(s, x, y, 0.045, h, color, None)

def notes(s, text):
    s.notes_slide.notes_text_frame.text = text

RULE = "실시간 빈티지 · 속보치 기준 · 전망주차 w[-19,-1] 평균 RMSE · 2018Q1~2025Q4(32개 분기) · 예측단위 개정(schema v2) 반영. 낮을수록 정확."

# ============ 1. 표지 ============
s = new_slide()
runs(s, 0.75, 0.55, 6.0, 0.35, [[("네이버클라우드 AX Forward Lab", True, GREY, 12)]])
runs(s, 0.75, 2.65, 11.5, 1.6, [
    [("GDP Nowcasting 협업 종합 리포트", True, NAVY, 33)],
    [("", False, INK, 8)],
    [("검증 총괄과 방법론 시사점, 그리고 공동 연구 제안", False, INK, 17)],
])
hline(s, 0.78, 4.35, 4.2, NAVY, 1.6)
runs(s, 0.78, 4.55, 11.0, 0.9, [
    [("두 달간 40여 개 모형 구성을 동일 규약으로 검증한 결과를 정리하고,", False, GREY, 12)],
    [("“사전학습 시계열 모델은 언제 가치를 더하는가”라는 공동 연구 방향을 제안드립니다.", False, GREY, 12)],
], sp=3)
runs(s, 0.78, 6.6, 6.0, 0.35, [[("2026. 8.  |  한국은행 디지털혁신실 협업", False, GREY, 11)]])
notes(s, "협업 기간(2026.6~8) 전체 검증 결과를 예측단위 개정(schema v2) 이후 기준으로 통합 정리한 자료. 모든 수치는 동일 채점 규약으로 재계산된 것이며, 개정 전 수치와의 혼용은 없음.")

# ============ 2. 리더보드 ============
s = new_slide()
header(s, "1  |  검증 총괄",
       "개정 이후 기준선은 XGBoost 단독 — 이를 넘는 구성은 주차별 결합 하나였습니다")
s.shapes.add_picture(f"{SC}/bokf_leader.png", Inches(0.5), Inches(1.6), width=Inches(7.5))
X = 8.45
sidebar(s, X, 1.75, 1.32)
runs(s, X + 0.18, 1.75, 4.15, 1.4, [
    [("개정 전 수치는 표에서 제외했습니다", True, NAVY, 11)],
    [("종전 최고로 보고되던 DFM+XGBoost 0.765는 예측", False, INK, 9.7)],
    [("단위 정합성 문제로 비교 불가 — 개정 후 같은 구성은", False, INK, 9.7)],
    [("0.787로, XGBoost 단독(0.750)보다 오히려 열세입니다.", False, INK, 9.7)],
], sp=2)
sidebar(s, X, 3.35, 1.32, GREEN)
runs(s, X + 0.18, 3.35, 4.15, 1.4, [
    [("유일하게 앞선 구성 (0.740, -1.3%)", True, GREEN, 11)],
    [("발표 19~14주 전에는 GBM·Chronos-2 예측의 평균,", False, INK, 9.7)],
    [("13주 전부터는 XGBoost 단독 — 주차별 교대입니다.", False, INK, 9.7)],
    [("다만 이 개선폭은 통계적 유의성에 미달합니다(3장).", False, INK, 9.7)],
], sp=2)
sidebar(s, X, 4.95, 1.1)
runs(s, X + 0.18, 4.95, 4.15, 1.2, [
    [("딥러닝 단독 최고는 0.808", True, NAVY, 11)],
    [("Chronos-2 + 일별신호. 기준선에는 미달하지만", False, INK, 9.7)],
    [("특정 조건에서 일관된 우위가 있습니다(3장).", False, INK, 9.7)],
], sp=2)
footnote(s, RULE + " 직접학습 신경망·학습형 결합 등 30여 개 구성은 전부 기준선 미달로 생략.")
notes(s, "리더보드는 schema v2(원 단위 예측) 기준으로 전면 재계산. 주차별 결합의 -1.3%는 DM 검정 유의성 미달(p=0.41)이며, 후보 구성 중 최선을 선택한 값이므로 '비열등+개선 방향' 수위로 보고. 통계적 우월 주장이 아님.")

# ============ 3. 시사점 ============
s = new_slide()
header(s, "2  |  시사점", "32개 분기 표본이 가르쳐준 것은 세 가지입니다")
items = [
    ("소표본에서는 단순한 고정 구조만 살아남았습니다", NAVY,
     [("결합 가중치 학습, 성과 기반 적응 가중, 잔차 보정 학습, 합성데이터 미세조정 — “데이터로부터 배우는”", False),
      ("개선책 네 계열을 모두 동일 규약으로 검증한 결과 전부 기준선 미달이었습니다. 살아남은 것은 사람이", False),
      ("고정한 단순 규칙(균등 평균, 주차별 교대)뿐입니다.", True)]),
    ("이 표본에서는 어떤 우월성도 통계적으로 입증되지 않습니다", NAVY,
     [("당사 검정력 분석으로는 32개 분기에서 6% 개선도 DM p=0.088에 그칩니다. 귀 기관의 MCS 검정에서도", False),
      ("비교 대상 12개 모형 전원이 5% Model Confidence Set에 잔류했습니다 — 양측 분석이 같은 결론입니다.", False),
      ("점 예측 정확도의 “순위 경쟁”은 이 표본에서 판정이 불가능한 게임입니다.", True)]),
    ("질문을 “누가 이기나”에서 “언제 기여하나”로 바꾸면 답이 있습니다", GREEN,
     [("모형 간 우열 대신 조건별 기여를 보면 — 해당 분기 지표가 발표되기 전인 조기 전망주차와 경기 반등기에서 사전학습", False),
      ("모델의 기여가 일관되게 관찰됩니다. 이 조건부 발견이 공동 연구 제안(6장)의 핵심 재료입니다.", False)]),
]
y = 1.75
for i, (title, c, body) in enumerate(items):
    sidebar(s, 0.6, y, 1.35, c)
    runs(s, 0.85, y, 11.7, 0.4, [[(f"{i+1}.  ", True, c, 13), (title, True, c, 13)]])
    runs(s, 1.15, y + 0.42, 11.3, 1.0, [[(t, b, INK, 10.5)] for t, b in body], sp=2)
    y += 1.72
footnote(s, RULE)
notes(s, "1의 '학습형 전패'는 우연이 아니라 구조적 결과 — 유효 표본 32개에서 추정하는 모든 추가 파라미터가 과적합 비용을 지불함. 2는 당사와 한은이 독립적으로 같은 결론에 도달했다는 점이 중요. 3이 본 리포트의 축.")

# ============ 4. 전망시계 진단 ============
s = new_slide()
header(s, "3  |  전망시계 진단", "해당 분기 지표가 나오기 전(조기 구간)에는 XGBoost의 우위가 사라집니다")
s.shapes.add_picture(f"{SC}/bokf_horizon.png", Inches(0.5), Inches(1.6), width=Inches(6.5))
X = 7.35
runs(s, X, 1.7, 5.4, 2.15, [
    [("조기 구간(발표 19~14주 전)은 해당 분기의 공식지표가 거의", False, INK, 10.5)],
    [("없는 시기입니다. 이때는 튜닝된 XGBoost의 우위가 사라지고,", False, INK, 10.5)],
    [("세계 지식과 일별 신호(주가·환율·심리 원지수)를 쓰는", False, INK, 10.5)],
    [("사전학습 모델이 근소하게 앞섭니다.", True, INK, 10.5)],
    [("", False, INK, 5)],
    [("경기 반등 분기 6개 평균에서도 같은 패턴입니다 —", False, INK, 10.5)],
    [("Chronos-2f 0.581 < XGBoost 0.717. 공식 지표의 발표", False, INK, 10.5)],
    [("시차를 일별 금융·심리 신호가 메우는 구간입니다.", False, INK, 10.5)],
], sp=2)
runs(s, X, 4.05, 5.4, 0.35, [[("리더보드 1위 구성 — 실제 결합 방식 (0.740)", True, GREEN, 12)]])
hline(s, X, 4.42, 5.4, GREEN, 1.0)
TL_X, TL_Y, TL_W = X + 0.05, 4.62, 5.25
w_early = TL_W * 6 / 19; w_rest = TL_W - w_early - 0.04
rect(s, TL_X, TL_Y, w_early, 0.4, WARM, None)
rect(s, TL_X + w_early + 0.04, TL_Y, w_rest, 0.4, BLUE, None)
runs(s, TL_X - 0.05, TL_Y + 0.06, w_early + 0.1, 0.3, [[("6주", True, WHITE, 10)]], align=PP_ALIGN.CENTER)
runs(s, TL_X + w_early + 0.04, TL_Y + 0.06, w_rest, 0.3, [[("13주", True, WHITE, 10)]], align=PP_ALIGN.CENTER)
runs(s, TL_X - 0.15, TL_Y + 0.46, w_early + 0.5, 0.6, [
    [("발표 19~14주 전", True, WARM, 9)],
    [("(GBM+Chronos-2f)÷2", False, INK, 9)]], align=PP_ALIGN.CENTER, sp=1)
runs(s, TL_X + w_early + 0.04, TL_Y + 0.46, w_rest, 0.6, [
    [("발표 13~1주 전", True, BLUE, 9)],
    [("XGBoost 단독 (현행 그대로)", False, INK, 9)]], align=PP_ALIGN.CENTER, sp=1)
runs(s, X, 5.85, 5.4, 0.7, [
    [("국면 판단은 없습니다 — 주차에 따른 고정 교대 규칙이며, 19주 중", False, GREY, 9.5)],
    [("13주는 현행 그대로입니다. 보고 수위는 “비열등 + 개선 방향”입니다.", False, GREY, 9.5)],
], sp=2)
footnote(s, RULE + " 반등 분기 6개: 2018Q1·2019Q2·2020Q3·2023Q1·2024Q3·2025Q2.")
notes(s, "구간 경계(-14주)는 이전 fan chart 검증에서 선정의된 것을 재사용(사후 선택 아님). 다만 조기 구간에 어떤 모델 조합을 쓸지는 후보 비교로 정했으므로 선택편의 가능성을 함께 명시하는 것이 정직한 보고.")

# ============ 5. ICL ============
s = new_slide()
header(s, "4  |  문맥학습(ICL) 계열", "성능 이전에, “공정한 채점” 자체가 성립하지 않습니다")
runs(s, 0.6, 1.42, 12.1, 0.3, [[("문맥학습(ICL) = 모델 파라미터를 학습하지 않고 과거 데이터를 문맥(프롬프트)으로 제시해 예측을 얻는 방식 — LLM 활용 계열 포함", False, GREY, 9.5)]])
X1, X2, W = 0.6, 7.0, 5.75
runs(s, X1, 1.95, W, 0.35, [[("LLM 기반 ICL의 구조적 한계", True, NAVY, 12.5)]])
hline(s, X1, 2.32, W, LINE, 1.0)
runs(s, X1, 2.48, W, 4.2, [
    [("학습데이터 오염 — 가장 결정적입니다. ", True, WARM, 10.5),
     ("과거 GDP 실적과 경제", False, INK, 10.5)],
    [("뉴스가 LLM 사전학습에 이미 들어 있어, “그 시점에 몰랐던 값”을", False, INK, 10.5)],
    [("모델이 알고 있을 가능성을 배제할 수 없습니다. 실시간 재현 채점", False, INK, 10.5)],
    [("자체가 성립하지 않아 리더보드에서 제외했습니다.", True, INK, 10.5)],
    [("당사 예비실험 결과도 ‘오염 상한’으로만 해석합니다.", False, GREY, 9.3)],
    [("", False, INK, 6)],
    [("재현성과 감사가능성. ", True, WARM, 10.5),
     ("동일 입력에도 출력이 변동하고 모델 버전", False, INK, 10.5)],
    [("변경 시 결과가 복원되지 않습니다 — 중앙은행 운영 요건(재현·", False, INK, 10.5)],
    [("감사·설명책임)과 정면으로 충돌합니다.", False, INK, 10.5)],
    [("", False, INK, 6)],
    [("수치 정밀도. ", True, WARM, 10.5),
     ("토큰 단위 수치 처리 특성상 소수점 회귀에 부적합.", False, INK, 10.5)],
], sp=2)
runs(s, X2, 1.95, W, 0.35, [[("직접학습 소형 신경망 이식판의 실측 (참고)", True, NAVY, 12.5)]])
hline(s, X2, 2.32, W, LINE, 1.0)
runs(s, X2, 2.48, W, 3.2, [
    [("직접학습 어텐션 LSTM 계열(과제특화 소형 모델)을 당사가 GDP", False, INK, 10.5)],
    [("과제용으로 재구현해 동일 규약으로 채점했습니다.", False, INK, 10.5)],
    [("", False, INK, 5)],
    [("· 단독 1.268 — 기준선(0.750)의 약 1.7배 오차", False, INK, 10.5)],
    [("· 경기 반등 분기 2.401 — 학습분포 밖 급변기 일반화 실패", False, INK, 10.5)],
    [("· XGBoost와 오차 상관 0.847 — 같은 패널을 소비해 결합", False, INK, 10.5)],
    [("  다양성이 없고, 모든 결합 배합에서 단조 열화", False, INK, 10.5)],
    [("· 원인: 월 140행 표본으로 다수 파라미터를 직접 학습", False, INK, 10.5)],
], sp=2.5)
rect(s, X2, 5.6, W, 0.95, BGGREY, None)
sidebar(s, X2, 5.6, 0.95, GREY)
runs(s, X2 + 0.2, 5.7, W - 0.4, 0.8, [
    [("유의: 본 결과는 당사 재구현 이식판에 대한 동일 규약 판정이며, 특정 원", False, GREY, 9.3)],
    [("연구나 타 과제(물가 등) 성능에 대한 판정이 아닙니다. 사전학습 FM 계열의", False, GREY, 9.3)],
    [("동일 프로토콜 비교는 별도 1p(부품 비교) 자료에 정리되어 있습니다.", False, GREY, 9.3)],
], sp=1.5)
footnote(s, RULE)
notes(s, "좌측: LLM ICL은 성능 이전에 평가 타당성(오염) 문제로 실시간 나우캐스팅 채점이 성립하지 않음 — 이것이 리더보드 제외 사유. 우측: 이식판 실측은 아키텍처 계열의 소표본 한계에 대한 증거이며, 특정 연구에 대한 평가가 아님을 명시적으로 구분.")

# ============ 6. 트랜스포머류 ============
s = new_slide()
header(s, "5  |  트랜스포머류(파운데이션 모델)", "점 예측의 대체재는 아니지만, 조건부 보완재로는 유효합니다")
X1, W1 = 0.6, 7.2
runs(s, X1, 1.6, W1, 0.35, [[("점 예측에 기여하기 어려운 네 가지 이유", True, NAVY, 12.5)]])
hline(s, X1, 1.97, W1, LINE, 1.0)
reasons = [
    ("소표본 정형 데이터는 GBDT의 영역입니다.",
     "월 140행·지표 34종 조건에서 튜닝된 트리 모델의 우위는 문헌 일반 결과(Grinsztajn et al., NeurIPS 2022)와 당사 실측(직접학습 신경망 전멸)이 일치합니다."),
    ("사전학습 지식은 한국 거시 구조가 아닙니다.",
     "파운데이션 모델이 배운 것은 추세·계절성 같은 시계열의 보편 문법입니다. 지표 간 인과나 한국 특수성은 없고, 공변량 소화력도 XGBoost(34종 비선형)에 미달합니다."),
    ("병목은 모델이 아니라 정보입니다.",
     "합성데이터로 추가 미세조정해도 예측이 거의 변하지 않았습니다(상관 0.990) — 입력 패널이 가진 정보 자체가 상한입니다."),
    ("분포 출력은 그대로 쓸 수 없습니다.",
     "원시 분위수의 실측 커버리지는 49%(명목 80%) — 과소산포로, 사후 보정이 필수입니다."),
]
y = 2.12
for i, (t, d) in enumerate(reasons):
    runs(s, X1, y, W1, 0.35, [[(f"{i+1}  ", True, WARM, 12), (t, True, INK, 11)]])
    runs(s, X1 + 0.28, y + 0.34, W1 - 0.3, 0.6, [[(d, False, INK, 9.8)]])
    y += 1.13
X2, W2 = 8.25, 4.5
rect(s, X2, 1.6, W2, 4.95, HL, None)
sidebar(s, X2, 1.6, 4.95, GREEN)
runs(s, X2 + 0.25, 1.78, W2 - 0.45, 4.6, [
    [("그럼에도 기여가 확인된 지점", True, GREEN, 12.5)],
    [("", False, INK, 5)],
    [("지표 공백 구간. ", True, INK, 10.5), ("발표 19~14주 전 — 해당 분기", False, INK, 10.5)],
    [("공식지표가 나오기 전의 공백을 세계 지식과", False, INK, 10.5)],
    [("일별 신호로 보완합니다(3장).", False, INK, 10.5)],
    [("", False, INK, 5)],
    [("경기 반등기. ", True, INK, 10.5), ("반등 6개 분기 평균 0.581로", False, INK, 10.5)],
    [("전 구성 중 최강 — 주가·환율·심리 원지수가", False, INK, 10.5)],
    [("공식 지표의 발표 시차를 메웁니다.", False, INK, 10.5)],
    [("", False, INK, 5)],
    [("불확실성 정량화 재료. ", True, INK, 10.5), ("보정(conformal)", False, INK, 10.5)],
    [("결합 시 예측구간 품질 개선을 확인 — 연구", False, INK, 10.5)],
    [("재료로 유효합니다.", False, INK, 10.5)],
    [("", False, INK, 6)],
    [("요약: 대체재가 아니라 조건부 보완재", True, GREEN, 11.5)],
], sp=2)
footnote(s, RULE)
notes(s, "핵심 메시지: '트랜스포머가 안 된다'가 아니라 '튜닝된 GBDT가 지배하는 조건과, 사전학습 지식이 기여하는 조건이 다르다'. 좌측 4개는 실측 근거가 각각 존재. 우측 조건부 기여가 공동 논문의 핵심 주장으로 연결.")

# ============ 7. 논문 ============
s = new_slide()
header(s, "6  |  공동 논문 제안", "“사전학습 시계열 모델은 언제 가치를 더하는가”를 묻는 논문을 제안드립니다")
runs(s, 0.6, 1.45, 12.1, 0.35, [
    [("가제: ", False, GREY, 10.5),
     ("When do time-series foundation models add value? A real-time evaluation of GDP nowcasting", True, NAVY, 11.5),
     ("   — 우월성 주장이 아닌 조건부 가치 규명 프레임", False, GREY, 10.5)],
])
X1, W1 = 0.6, 6.5
runs(s, X1, 2.0, W1, 0.35, [[("기여 다섯 가지 — 실측은 모두 완료되어 있습니다", True, NAVY, 12)]])
hline(s, X1, 2.37, W1, LINE, 1.0)
contribs = [
    ("2세대 TSFM의 실시간 빈티지 평가 최초 사례", "주 단위 19회 예측 그리드 — 기존 문헌은 대부분 사후 확정 데이터"),
    ("전망시계 조건부 보완성", "정보 빈곤 구간·반등기에서만 기여 — 메커니즘 있는 조건부 발견"),
    ("무개정 일별 신호의 식별 깨끗한 효과", "가격 데이터는 개정이 없어 look-ahead 논란 원천 차단 (0.836→0.808)"),
    ("TSFM 분위수 과소산포 실측과 처방", "커버리지 49% → 컨포멀 보정으로 복원"),
    ("소표본 결합 퍼즐 + 검정력 분석", "학습형 전패·고정 구조 생존, MCS 전원 잔류 — ‘평가 축 재설계’ 논의로 승격"),
]
y = 2.52
for i, (t, d) in enumerate(contribs):
    runs(s, X1, y, W1, 0.3, [[(f"{i+1}. ", True, GREEN, 10.5), (t, True, INK, 10.5)]])
    runs(s, X1 + 0.25, y + 0.29, W1 - 0.25, 0.3, [[(d, False, GREY, 9)]])
    y += 0.66
runs(s, X1, y + 0.1, W1, 0.7, [
    [("역할 분담(안)  ", True, NAVY, 10.5), ("한은: 실시간 빈티지·도메인 검증·정책 함의  /  당사: 모델·평가 인프라·재현 패키지", False, INK, 10)],
    [("발표 경로(안)  ", True, NAVY, 10.5), ("한은 Working Paper 선행 공개 → International Journal of Forecasting 투고", False, INK, 10)],
], sp=3)
X2, W2 = 7.55, 5.2
runs(s, X2, 2.0, W2, 0.35, [[("표본 한계는 Future Work가 아니라 설계로 풉니다", True, NAVY, 12)]])
hline(s, X2, 2.37, W2, LINE, 1.0)
runs(s, X2, 2.52, W2, 4.0, [
    [("다국가 복제. ", True, INK, 10.5), ("미국 ALFRED(세인트루이스 연준 실시간", False, INK, 10.5)],
    [("DB, 공개)·유로존 ECB 실시간 DB로 동일 프레임을 재현해,", False, INK, 10.5)],
    [("‘단일국 32분기’라는 예상 반론에 선제 대응합니다.", False, INK, 10.5)],
    [("", False, INK, 6)],
    [("사전등록 라이브 추적. ", True, INK, 10.5), ("주차별 결합 구성을 고정·공개한", False, INK, 10.5)],
    [("뒤 향후 분기를 실전 축적 — 선택편의 없는 표본외 증거를", False, INK, 10.5)],
    [("만듭니다.", False, INK, 10.5)],
    [("", False, INK, 6)],
    [("유의성 논의는 본문에서 정면 처리. ", True, INK, 10.5), ("검정력 분석과", False, INK, 10.5)],
    [("MCS를 ‘평가 방법론 기여’로 배치합니다 — 약점 고백이", False, INK, 10.5)],
    [("아니라 논문 주장의 일부입니다.", False, INK, 10.5)],
], sp=2)
footnote(s, "데이터 공개 방식(코드 공개 + 빈티지 접근 절차)은 초기 합의 필요. 기준선 개정 경위의 서술 수위는 공동저자 협의 사항.")
notes(s, "제목·프레임의 핵심은 '이긴다'가 아니라 '언제 기여하는가'. 이 구도에서는 유의성 미달이 약점이 아니라 5번 기여(평가 방법론)의 일부가 됨. 다국가 복제는 리뷰어 1순위 예상 공격(단일국 32분기)에 대한 선제 설계.")

# ============ 8. 요약 ============
s = new_slide()
header(s, "7  |  요약", "네 가지로 정리드리고, 논문 구도 협의를 요청드립니다")
rows = [
    ("점 예측", "XGBoost 단독(개정 후)이 기준선 — 인정합니다. 주차별 결합(조기 6주 = GBM·Chronos-2 평균 → 이후 XGBoost, 0.740)은 비열등+개선 방향의 참고 옵션입니다.", "필요 시 병행 산출"),
    ("문맥학습(ICL)", "LLM 계열은 오염 문제로 실시간 채점 자체가 성립하지 않고, 직접학습 이식판은 소표본 한계를 실측했습니다.", "채택 비권고"),
    ("트랜스포머류", "점 예측 대체재가 아니라 지표 공백 구간(조기 주차)·반등기의 조건부 보완재이며, 불확실성 정량화의 재료입니다.", "논문 핵심 재료"),
    ("공동 논문", "“언제 가치를 더하는가” 프레임으로 기여 5건의 실측이 완료되어 있고, 다국가 복제 설계를 포함합니다.", "구도 협의 요청"),
]
y = 1.75
for a, bb, c in rows:
    hline(s, 0.6, y - 0.12, 12.13, LINE, 0.75)
    runs(s, 0.65, y, 2.0, 0.4, [[(a, True, NAVY, 12.5)]])
    runs(s, 2.85, y, 7.6, 0.85, [[(bb, False, INK, 10.5)]])
    runs(s, 10.7, y + 0.12, 2.0, 0.4, [[(c, True, GREEN, 11)]], align=PP_ALIGN.RIGHT)
    y += 1.12
hline(s, 0.6, y - 0.12, 12.13, LINE, 0.75)
runs(s, 0.6, y + 0.25, 12.1, 0.4, [
    [("모든 실험은 재현 가능한 스크립트로 관리하고 있으며, 논문 착수 결정 시 재현 패키지 형태로 정리해 공유하겠습니다.", False, GREY, 10.5)]])
footnote(s, RULE)
notes(s, "다음 회의 의사결정 요청 사항: ① 공동 논문 구도(제목 프레임·역할 분담·발표 경로) 협의 착수 여부 ② 주차별 결합 병행 산출 여부(선택). 나머지는 보고 사항.")

# ============ 9. 별첨: 시계열 FM 부품 비교 ============
s = new_slide()
header(s, "별첨  |  시계열 FM 부품 비교",
       "부품으로서의 시계열 FM(ICL) 비교 — 조기 슬롯 자격은 Chronos-2만 충족")
runs(s, 0.6, 1.4, 12.1, 0.6, [
    [("실험 설계  ", True, NAVY, 10),
     ("검증된 주차별 결합의 조기 슬롯에 시계열 FM을 부품으로 장착해 비교. 두 모델 모두 파라미터 갱신 없이 과거 데이터를 문맥으로 사용(ICL·zero-shot), "
      "입력·과제·평가 완전 동일 — DFM 스냅샷 + 공변량 10종 + 빠른신호 4종, 분기말 외삽, 32분기 실시간 그리드.", False, INK, 10)],
    [("비교 대상  ", True, NAVY, 9.3),
     ("Chronos-2(아마존, 120M) vs BISTRO(BIS WP 1337, Moirai-base 91M을 BIS 거시 4,925계열로 미세조정 — 공개 체크포인트) · 참고: 직접학습 LSTM · LLM 기반 ICL은 학습데이터 오염으로 채점 제외.", False, GREY, 9.3)],
])
s.shapes.add_picture(f"{SC}/fm_parts.png", Inches(0.5), Inches(2.15), width=Inches(7.3))
runs(s, 0.6, 5.72, 7.1, 0.3, [[("결과 요약", True, NAVY, 11)]])
hline(s, 0.6, 6.02, 7.1, LINE, 0.75)
runs(s, 0.6, 6.1, 7.1, 0.85, [
    [("· 단독 성능: Chronos-2f 0.808 < BISTRO 0.871 ≈ 기반 Moirai-small 0.873 — BIS 거시 미세조정의 효과는 소폭(조기 구간 1.021→1.005)", False, INK, 9)],
    [("· 슬롯 장착: Chronos-2f만 개선(0.750→0.740). BISTRO는 동률(0.750), 두 FM 병용(0.744)도 C2f 단독에 미달", False, INK, 9)],
    [("· 참고: 직접학습 LSTM은 동일 입력에도 1.019(조기 1.286) — 소표본 직접학습 한계 재확인", False, INK, 9)],
], sp=2)
AX, AW = 8.15, 4.6
sidebar(s, AX, 2.15, 1.45, GREEN)
runs(s, AX + 0.2, 2.15, AW - 0.2, 1.5, [
    [("슬롯 자격 — 조기 구간에서 GBM과 대등할 것", True, NAVY, 11)],
    [("· 조기 구간 단독: Chronos-2f 0.924 ≈ GBM 0.924 (충족) /", False, INK, 9.3)],
    [("  BISTRO 1.005 · LSTM 1.286 (미충족)", False, INK, 9.3)],
    [("· 슬롯 결과가 자격 판정과 일치: C2f 0.740(개선) / BISTRO", False, INK, 9.3)],
    [("  0.750(동률) / 병용 0.744 — C2f 단독 미달 (오차 상관 0.926)", False, INK, 9.3)],
], sp=2)
sidebar(s, AX, 3.75, 2.45, NAVY)
runs(s, AX + 0.2, 3.75, AW - 0.2, 2.5, [
    [("성능 차이의 원인 — 공변량 절제 실험", True, NAVY, 11)],
    [("· 같은 빠른신호(주가·환율) 추가 시: ", False, INK, 8.8),
     ("C2 -3.4%", True, GREEN, 8.8), (" (0.836→0.808)", False, GREY, 8.3)],
    [("  vs ", False, INK, 8.8), ("BISTRO +0.1% 무이득", True, WARM, 8.8), (" (0.870→0.871)", False, GREY, 8.3)],
    [("· 공식지표 10종 추가: BISTRO -0.5% — 미미", False, INK, 8.8)],
    [("· 원인은 사전학습의 채점 범위 차이:", False, INK, 8.8)],
    [("  - Moirai/BISTRO: (공변량, 타깃)과거 → ", False, INK, 8.8),
     ("(공변량, 타깃)미래 전부 채점", True, INK, 8.8)],
    [("    = 공변량도 예측 대상(이웃). 거시 미세조정 후에도 구조 동일", False, INK, 8.8)],
    [("  - Chronos-2: (공변량, 타깃)과거 → ", False, INK, 8.8),
     ("타깃만 채점", True, INK, 8.8)],
    [("    = 공변량은 힌트. 힌트 없이는 손실을 줄일 수 없는 사례로 훈련", False, INK, 8.8)],
    [("· 체급 유사(91M vs 120M) → 격차 원인은 체급이 아닌 훈련 방식", False, INK, 8.8)],
], sp=1.6)
rect(s, AX, 6.35, AW, 0.7, BGGREY, None)
sidebar(s, AX, 6.35, 0.7, GREY)
runs(s, AX + 0.2, 6.4, AW - 0.4, 0.6, [
    [("유의  개선폭은 모두 통계적 유의성 미달(DM p>0.18) — 비열등·개선 방향 수위.", False, GREY, 8.3)],
    [("BISTRO는 공개 체크포인트를 논문 표준 사용법(zero-shot) 그대로 적용, 한국 GDP", False, GREY, 8.3)],
    [("과제용 추가 미세조정 없음. 저자: Koyuncu·Kwon·Lombardi·Perez-Cruz·Shin.", False, GREY, 8.3)],
], sp=1.3)
footnote(s, RULE + " 슬롯 교체 = 조기(발표 19~14주 전)=(GBM+부품)÷2, 이후 XGBoost. BISTRO 출처: github.com/bis-med-it/bistro (Apache 2.0).")
notes(s, "별첨 취지: 본문 결론(트랜스포머류 = 조건부 보완재)을 바꾸지 않는 보강 자료. 한은 측의 'Chronos-2 vs BISTRO 차이' 질문에 대응. 회의에서 다룰지는 현장 판단 — 본문 7장까지로 마쳐도 완결됨. BISTRO 관련 서술은 사실·수치 중심으로 유지(공개 체크포인트, 논문 표준 사용법 명시).")

out = "/Users/user/vibe/bistro-lstm/docs/GDP_Nowcasting_종합리포트_한은_2026-08-10.pptx"
prs.save(out); print("saved:", out)
