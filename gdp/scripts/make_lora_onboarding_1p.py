# -*- coding: utf-8 -*-
"""인턴 온보딩 1p — LoRA 적응 학습의 데이터·기간·방법 (문장형, 나눔스퀘어)"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

F = "NanumSquare"
SCALE = 1.15
INK = RGBColor(0x26,0x26,0x26); GREY = RGBColor(0x6B,0x6B,0x6B)
LINE = RGBColor(0xC9,0xC9,0xC9); NAVY = RGBColor(0x24,0x34,0x47)
GREEN = RGBColor(0x00,0x8A,0x3E); BGGREY = RGBColor(0xF5,0xF5,0xF4)
WARM = RGBColor(0xB0,0x53,0x2F)

SC = "/private/tmp/claude-502/-Users-user-vibe-bistro-lstm/031ca4b9-003d-4cb2-afaa-06353432dc9b/scratchpad"
prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
s = prs.slides.add_slide(prs.slide_layouts[6])

def runs(x, y, w, h, lines, size=10, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp=2):
    tf = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.02); tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01); tf.margin_bottom = Inches(0.01)
    for i, rs in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp)
        for t, b, c, z in rs:
            r = p.add_run(); r.text = t
            r.font.name = F; r.font.size = Pt(round((z if z else size) * SCALE, 1))
            r.font.bold = b; r.font.color.rgb = c if c else INK
            rp = r._r.get_or_add_rPr(); rp.append(rp.makeelement(qn('a:ea'), {'typeface': F}))
    return tf

def hline(x, y, w, color=LINE, pt=0.75):
    ln = s.shapes.add_connector(1, Inches(x), Inches(y), Inches(x+w), Inches(y))
    ln.line.color.rgb = color; ln.line.width = Pt(pt); return ln

def sidebar(x, y, h, color=NAVY):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.045), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = color; shp.line.fill.background(); shp.shadow.inherit = False

# 헤더
runs(0.6, 0.3, 11.5, 0.3, [[("온보딩 노트  |  시계열 Track1 · 적응학습(LoRA)", True, GREY, 10)]])
runs(0.6, 0.58, 12.1, 0.6, [[("LoRA 적응은 무엇을 보고 배우나 — 데이터, 기간, 학습 방법", True, NAVY, 16)]])
hline(0.6, 1.24, 12.13, NAVY, 1.2)

# 좌: 다이어그램
s.shapes.add_picture(f"{SC}/lora_dataset.png", Inches(0.45), Inches(1.5), width=Inches(6.6))

# 우: 문장형 설명
X, W = 7.35, 5.4
runs(X, 1.5, W, 1.5, [
    [("데이터셋은 어떻게 생겼나", True, NAVY, 12)],
    [("표 형태의 큰 파일이 아니라, 분기마다 하나씩 만든 짧은 시계열 묶음입니다.", False, INK, 9.8)],
    [("한 시리즈는 그 분기 시점에 실제로 알 수 있었던 월별 데이터(GDP 경로,", False, INK, 9.8)],
    [("공식지표 10종, 주가·환율 같은 빠른신호 4종, 그리고 각 달이 실측인지", False, INK, 9.8)],
    [("추정인지 표시하는 플래그)로 이루어져 있고, 분기말 달의 값만 나중에", False, INK, 9.8)],
    [("발표된 실제 속보치로 바꿔 넣습니다. 이 마지막 값이 정답 역할을 합니다.", False, INK, 9.8)],
], sp=2)
runs(X, 3.15, W, 1.15, [
    [("기간은 어떻게 되나", True, NAVY, 12)],
    [("각 시리즈의 과거 구간은 최대 128개월(약 11년)입니다. 적응은 해마다", False, INK, 9.8)],
    [("한 번, 그 해 이전에 속보가 발표 완료된 분기들만 모아서 합니다. 예를", False, INK, 9.8)],
    [("들어 2024년용 적응에는 23개 분기가 들어가고, 평가는 2021~2025년의", False, INK, 9.8)],
    [("분기들로 합니다. 학습에 쓴 분기는 채점에 절대 들어가지 않습니다.", False, INK, 9.8)],
], sp=2)
runs(X, 4.55, W, 1.2, [
    [("학습은 어떻게 하나", True, NAVY, 12)],
    [("모델 본체(약 1억 파라미터)는 손대지 않고 얼려 둡니다. 대신 주의", False, INK, 9.8)],
    [("(attention) 층 옆에 전체의 0.6%쯤 되는 작은 우회 회로(LoRA, rank 8)를", False, INK, 9.8)],
    [("붙이고 그것만 학습합니다. 시리즈들을 보여주며 \"경로의 마지막 6개월을", False, INK, 9.8)],
    [("맞혀 보라\"를 500번 반복하는 것이 전부입니다. 한 번 적응에 수 분이 걸립니다.", False, INK, 9.8)],
], sp=2)
sidebar(X, 5.95, 1.15, WARM)
runs(X + 0.2, 5.95, W - 0.2, 1.15, [
    [("과적합은 안 생기나", True, WARM, 12)],
    [("생깁니다. 적응 재료가 7분기도 안 되던 시기에는 적응 안 한 모델보다", False, INK, 9.8)],
    [("오히려 나빠지는 것을 직접 확인했습니다. 그래서 \"적응 이력이 12분기", False, INK, 9.8)],
    [("넘게 쌓인 뒤에만 쓴다\"는 규칙과 seed 3회 평균이 설계에 들어 있습니다.", False, INK, 9.8)],
], sp=2)

# 각주
runs(0.6, 7.08, 11.6, 0.3, [
    [("관련 코드: phase_b_c2_lora.py(Chronos-2) · phase_b_bistro_lora.py(BISTRO) — gdp/scripts. 플랜 문서: GDP_LoRA적응_실험플랜_2026-08-26.md", False, GREY, 7.5)]])

out = "/Users/user/vibe/bistro-lstm/docs/온보딩_LoRA학습구조_2026-09-02.pptx"
prs.save(out); print("saved:", out)
