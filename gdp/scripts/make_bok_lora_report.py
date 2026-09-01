# -*- coding: utf-8 -*-
"""한은 미팅(목) 경과 보고 — 적응학습 실증 결과와 운영 접목 (6장, 나눔스퀘어)"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

F = "NanumSquare"
SCALE = 1.15
INK = RGBColor(0x26,0x26,0x26); GREY = RGBColor(0x6B,0x6B,0x6B)
LINE = RGBColor(0xC9,0xC9,0xC9); NAVY = RGBColor(0x24,0x34,0x47)
GREEN = RGBColor(0x00,0x8A,0x3E); BGGREY = RGBColor(0xF5,0xF5,0xF4)
WHITE = RGBColor(0xFF,0xFF,0xFF); HL = RGBColor(0xEE,0xF5,0xF0)
WARM = RGBColor(0xB0,0x53,0x2F); BLUE = RGBColor(0x1C,0x5C,0xAB)

SC = "/private/tmp/claude-502/-Users-user-vibe-bistro-lstm/031ca4b9-003d-4cb2-afaa-06353432dc9b/scratchpad"
prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
PAGE = [0]
def new_slide():
    PAGE[0] += 1
    return prs.slides.add_slide(prs.slide_layouts[6])

def runs(s, x, y, w, h, lines, size=10, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp=2):
    tf = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.02); tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01); tf.margin_bottom = Inches(0.01)
    for i, rs in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp)
        for t, b, c, z in rs:
            r = p.add_run(); r.text = t
            r.font.name = F; r.font.size = Pt(round((z if z else size) * SCALE, 1))
            r.font.bold = b; r.font.color.rgb = c if c else INK
            rp = r._r.get_or_add_rPr(); rp.append(rp.makeelement(qn('a:ea'), {'typeface': F}))
    return tf

def hline(s, x, y, w, color=LINE, pt=0.75):
    ln = s.shapes.add_connector(1, Inches(x), Inches(y), Inches(x+w), Inches(y))
    ln.line.color.rgb = color; ln.line.width = Pt(pt); return ln

def rect(s, x, y, w, h, fill=None, line=None):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False; return shp

def header(s, tag, msg):
    runs(s, 0.6, 0.3, 11.5, 0.3, [[(tag, True, GREY, 10)]])
    runs(s, 0.6, 0.58, 12.1, 0.6, [[(msg, True, NAVY, 16)]])
    hline(s, 0.6, 1.24, 12.13, NAVY, 1.2)

def footnote(s, text):
    runs(s, 0.6, 7.08, 11.4, 0.3, [[("주: " + text, False, GREY, 7.5)]])
    runs(s, 12.85, 7.08, 0.4, 0.25, [[(str(PAGE[0]), False, GREY, 8.5)]], align=PP_ALIGN.RIGHT)

def sidebar(s, x, y, h, color=NAVY):
    rect(s, x, y, 0.045, h, color, None)

def notes(s, text):
    s.notes_slide.notes_text_frame.text = text

RULE = "실시간 빈티지 · 속보치 기준 · 주차별 RMSE(분기 평균) · 예측단위 개정(schema v2) 반영 · 낮을수록 정확"

# ============ 1. 표지 ============
s = new_slide()
runs(s, 0.75, 0.55, 8.0, 0.35, [[("네이버클라우드 AX Forward Lab", True, GREY, 11)]])
runs(s, 0.75, 2.5, 11.8, 1.7, [
    [("적응학습(LoRA) 실증 결과와 운영 접목 방안", True, NAVY, 28)],
    [("", False, INK, 8)],
    [("귀 리포트(8/25) 제안에 대한 실측 답변 — 협업 경과 보고", False, INK, 15)],
])
hline(s, 0.78, 4.3, 4.2, NAVY, 1.6)
runs(s, 0.78, 4.5, 11.2, 1.0, [
    [("두 파운데이션 모델(BISTRO·Chronos-2)에 동일 레시피의 LoRA 적응을 적용해 seed 3회로 검증했고,", False, GREY, 11.5)],
    [("적응 부품을 운영 구성(주차별 결합)에 접목한 결과까지 정리했습니다.", False, GREY, 11.5)],
], sp=3)
runs(s, 0.78, 6.55, 9.0, 0.35, [[("2026. 9. 4. 회의용  |  시계열 Track1 (거시경제 예측)  |  네이버클라우드 김용민", False, GREY, 10.5)]])
notes(s, "목요일 회의 안건: ①적응학습 실증 결과 공유 ②운영 구성(임계 하이브리드) 협의 ③역할 분담·재현 패키지·논문 장 구성. 이 자료는 8/25 이창훈 과장 리포트의 세 논의사항에 대한 실측 기반 답변.")

# ============ 2. 경과 요약 ============
s = new_slide()
header(s, "1  |  경과 요약", "귀 리포트(8/25) 수신 이후 2주간 수행한 실험입니다")
items = [
    ("8/26", "빈티지 흡수 진단", "귀측 진단(BISTRO의 주차별 오차 평탄성)을 당사 규약에서 재검증 — 평탄성은 모델 고유 특성이며, 같은 조건의 Chronos-2는 흡수율 2배로 우하향함을 확인", NAVY),
    ("8/26~27", "LoRA 적응 실증", "귀측 제안(빈티지 경로 기반 제한적 추가학습)을 LoRA로 구현 — BISTRO·Chronos-2 모두에서 개선 확인, seed 3회 반복으로 재현성 검증", GREEN),
    ("8/28", "운영 구성 접목", "적응 부품을 주차별 결합의 조기 슬롯에 장착 — 32분기 전체 검증에서 적응 이력 임계 규칙 포함 시 기존 최고 구성을 경신(0.740 → 0.733)", GREEN),
    ("병행", "품질 관리", "seed 미전파 결함 발견·수정, 실패 변형(주차 절단 경로 학습) 기각 기록, 전 실험 스크립트·결과 커밋 보존", GREY),
]
y = 1.55
for when, title, desc, c in items:
    sidebar(s, 0.6, y, 1.05, c)
    runs(s, 0.85, y, 1.35, 0.4, [[(when, True, c, 11.5)]])
    runs(s, 2.3, y, 2.3, 0.4, [[(title, True, NAVY, 11.5)]])
    runs(s, 4.7, y, 7.95, 0.95, [[(desc, False, INK, 10)]])
    y += 1.28
footnote(s, RULE + ". 상세 수치·스크립트는 회의 시 재현 패키지와 함께 공유 가능.")
notes(s, "핵심 메시지: 귀측의 세 논의사항(기준축 유지 여부 / 협업 범위 확대 / 역할 분담)에 말이 아니라 실측으로 답해왔음. 4행(품질 관리)은 신뢰 서사 — 실패와 결함도 기록으로 남긴다.")

# ============ 3. 흡수 진단 ============
s = new_slide()
header(s, "2  |  빈티지 흡수 진단", "귀측 진단이 정확했습니다 — 평탄성은 모델 고유 특성이며, 출발 체크포인트에 따라 갈립니다")
s.shapes.add_picture(f"{SC}/weekly_absorb.png", Inches(0.5), Inches(1.6), width=Inches(7.3))
X = 8.1
sidebar(s, X, 1.7, 1.7, WARM)
runs(s, X + 0.2, 1.7, 4.6, 1.75, [
    [("진단 재현", True, NAVY, 11.5)],
    [("· 일별신호까지 동일하게 준 조건에서도 BISTRO는", False, INK, 9.5)],
    [("  주차 경과에 둔감 (흡수율 -4.7%)", False, INK, 9.5)],
    [("· 기반 Moirai-small도 -7.0% — 계열 공통 특성", False, INK, 9.5)],
    [("· 귀 리포트의 '빈티지 정보 미반영' 진단과 일치", False, INK, 9.5)],
], sp=2)
sidebar(s, X, 3.6, 1.7, GREEN)
runs(s, X + 0.2, 3.6, 4.6, 1.75, [
    [("같은 조건의 Chronos-2는 우하향", True, NAVY, 11.5)],
    [("· 흡수율 -9.3% (BISTRO의 약 2배)", False, INK, 9.5)],
    [("· 원인: 사전학습의 채점 범위 차이 — 전 변량을 함께", False, INK, 9.5)],
    [("  예측하도록 훈련 vs 타깃만 채점(공변량=힌트)", False, INK, 9.5)],
    [("· 절제 실험과 정합: 같은 일별신호에서 C2 -3.4% vs", False, INK, 9.5)],
    [("  BISTRO 무이득", False, INK, 9.5)],
], sp=2)
runs(s, X, 5.5, 4.6, 1.3, [
    [("→ 이 격차가 적응학습의 출발점 질문이 됩니다:", False, INK, 10)],
    [("\"경량 적응이 흡수 능력을 만들어줄 수 있는가\" (다음 장)", True, NAVY, 10.5)],
], sp=2)
footnote(s, RULE + ". 흡수율 = 초반 6주 대비 다음 6주 RMSE 변화(w=-19~-8, TSFM 자체 예측 구간). 절제 = 공변량 단계적 추가로 기여 분리.")
notes(s, "톤: '귀측 진단이 정확했다'로 시작 — 진단 재현이 협업 예의이자 신뢰 형성. BISTRO 관련 서술은 사실·수치만, 항상 '공개 체크포인트·동일 규약' 프레임.")

# ============ 4. LoRA 실증 ============
s = new_slide()
header(s, "3  |  적응학습(LoRA) 실증", "제안하신 방향이 유효합니다 — 두 모델 모두 개선, BISTRO는 평탄성도 부분 해소")
s.shapes.add_picture(f"{SC}/lora_share.png", Inches(0.5), Inches(1.65), width=Inches(7.3))
X = 8.1
sidebar(s, X, 1.7, 1.9, GREEN)
runs(s, X + 0.2, 1.7, 4.6, 1.95, [
    [("결과 (3-seed 예측 평균)", True, NAVY, 11.5)],
    [("· BISTRO  0.596 → 0.560 (-6.0%), 흡수 역행", False, INK, 9.5)],
    [("  +10.1% → +3.2%로 완화", False, INK, 9.5)],
    [("· Chronos-2f  0.619 → 0.564 (-8.9%)", False, INK, 9.5)],
    [("· 귀측 GDP head 실험(출력층 교체)과의 대조 —", False, INK, 9.5)],
    [("  주의층 0.6%(rank 8) 적응이 관건이었음", True, INK, 9.5)],
], sp=2)
sidebar(s, X, 3.8, 1.55, WARM)
runs(s, X + 0.2, 3.8, 4.6, 1.6, [
    [("유의", True, NAVY, 11.5)],
    [("· BISTRO 쪽 seed 편차 큼(0.549~0.602) — 단일 실행", False, INK, 9.5)],
    [("  불신, 3-seed 예측 평균을 대표값으로", False, INK, 9.5)],
    [("· 개선폭의 통계적 유의성은 표본 한계로 미검정", False, INK, 9.5)],
    [("· 사전학습-평가기간 중복은 귀 리포트와 동일하게", False, INK, 9.5)],
    [("  회고적 진단 수위로 해석", False, INK, 9.5)],
], sp=2)
runs(s, X, 5.55, 4.6, 1.2, [
    [("학습 설계 = 귀측 제안 그대로: 분기별 빈티지 경로(당시", False, GREY, 9)],
    [("가용 정보 + 관측 플래그, 라벨=실제 속보치), 연 1회 fresh", False, GREY, 9)],
    [("재적응, release-safe, walk-forward.", False, GREY, 9)],
], sp=1.5)
footnote(s, "평가창 2021~2025 (적응 데이터 확보 구간, 20분기). BISTRO LoRA는 patch=8 고정 — zero-shot 대조군도 동일 patch(공정 비교).")
notes(s, "정치적 핵심: '제안하신 방향이 유효합니다' — 이과장 제안의 실증 확인이라는 프레임. GDP head(실패)와의 대조는 사실 서술로만, 그들의 실험을 존중하는 톤으로 구두 보완.")

# ============ 5. 운영 접목 ============
s = new_slide()
header(s, "4  |  운영 구성 접목", "적응 부품 + 이력 임계 규칙으로 32분기 기준 기존 최고 구성을 경신했습니다")
runs(s, 0.6, 1.5, 6.2, 2.4, [
    [("구성 (주차별 결합의 조기 슬롯 부품 교체)", True, NAVY, 12)],
    [("· 기본 = XGBoost 단독 (개정 후 기준선 0.750)", False, INK, 10)],
    [("· 조기 6주(지표 공백기) = FM 결합으로 교대", False, INK, 10)],
    [("   - 적응 이력 12분기 미만 → zero-shot 부품 (GBM+C2)/2", False, INK, 10)],
    [("   - 12분기 이상 → LoRA 적응 부품 (C2-LoRA+BISTRO-LoRA)/2", False, INK, 10)],
    [("· 임계 12분기는 실험 전 플랜에 명시했던 사전 등록값", False, INK, 10)],
])
t_rows = [("구성", "32분기 전체", "분기별 승수"),
          ("XGBoost 단독 (기준선)", "0.750", "—"),
          ("+ 조기 6주 zero-shot 결합 (기존 공유안)", "0.740", "16 / 32"),
          ("+ 적응 이력 임계 규칙 (신규)", "0.733", "18 / 32")]
ty = 4.05
tbl = s.shapes.add_table(4, 3, Inches(0.6), Inches(ty), Inches(6.2), Inches(1.7)).table
tbl.columns[0].width = Inches(3.6); tbl.columns[1].width = Inches(1.3); tbl.columns[2].width = Inches(1.3)
for r in range(4):
    for c in range(3):
        cell = tbl.cell(r, c)
        cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.04)
        cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY if r == 0 else (HL if r == 3 else WHITE)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
        run = p.add_run(); run.text = t_rows[r][c]
        run.font.name = F; run.font.size = Pt(round(9.5 * SCALE, 1)); run.font.bold = (r in (0, 3))
        run.font.color.rgb = WHITE if r == 0 else INK
        rp = run._r.get_or_add_rPr(); rp.append(rp.makeelement(qn('a:ea'), {'typeface': F}))
X = 7.2
sidebar(s, X, 1.5, 2.1, WARM)
runs(s, X + 0.2, 1.5, 5.5, 2.15, [
    [("임계 규칙이 필요한 이유 (실측)", True, NAVY, 11.5)],
    [("· 적응 데이터 3~7분기 시기(2019~20, 코로나 포함)에", False, INK, 9.5)],
    [("  LoRA 부품을 쓰면 zero-shot보다 오히려 악화", False, INK, 9.5)],
    [("  (조기 구간 1.71 vs 1.50) — 소형 적응의 위험 확인", False, INK, 9.5)],
    [("· 충분한 이력(12분기+)에서만 적응 부품이 자격을 가짐", False, INK, 9.5)],
    [("· \"언제부터 적응 모델을 신뢰할 수 있는가\"의 기준을", False, INK, 9.5)],
    [("  실측으로 제시 — 논문 재료로도 유효", False, INK, 9.5)],
], sp=2)
sidebar(s, X, 3.85, 1.35, NAVY)
runs(s, X + 0.2, 3.85, 5.5, 1.4, [
    [("수위 (일관 유지)", True, NAVY, 11.5)],
    [("· -2.3%는 DM p=0.185 — 통계적 유의성 미달,", False, INK, 9.5)],
    [("  \"비열등 + 개선 방향\"으로만 제시", False, INK, 9.5)],
    [("· 단, 검증한 전 구성 중 가장 낮은 p·최다 승수", False, INK, 9.5)],
    [("· 운영 시 seed 3개 × 2모델 러너 유지 비용 존재", False, INK, 9.5)],
], sp=2)
footnote(s, RULE + " · 전망주차 w[-19,-1] · 2018Q1~2025Q4(32개 분기). 임계 미충족 구간(2018~2020)은 zero-shot 부품 사용.")
notes(s, "이 장이 회의의 실질 협의 대상 — 병행 산출 여부. 임계 12분기의 사전 등록 근거(플랜 문서 §3)를 회의에서 바로 제시할 수 있게 준비. 2024년 소폭 열세(0.652 vs 0.628)도 질문 시 숨기지 않고 공유.")

# ============ 6. 제안 ============
s = new_slide()
header(s, "5  |  협의 안건", "역할 분담과 다음 단계를 정하고자 합니다")
agenda = [
    ("A1", "적응학습 역할 분담", "당사: Chronos-2 트랙 · 적응 안정화(seed 앙상블) · 흡수 진단 하네스 / 귀측: BISTRO 트랙 · 공표달력(release mask) 정밀 정의 · raw vintage 입력 — 사전 등록 설정 합의 후 각자 실행, 동일 규약 교차 채점", GREEN),
    ("A2", "재현 패키지 교차 검증", "당사 실험 스크립트 일체(LoRA 러너·채점 하네스 포함)를 재현 가능 형태로 공유 — 귀측 인프라에서 독립 재현", NAVY),
    ("A3", "운영 구성 병행 산출", "임계 하이브리드(0.733)를 현행과 나란히 참고 지표로 산출·축적할지 여부 — 교체가 아닌 병행 검증", NAVY),
    ("A4", "공동 논문 적응학습 장(章)", "\"언제 가치를 더하는가\" 프레임에 적응학습 결과 추가 — zero-shot 진단 → 출력층 → LoRA → 임계 규칙의 완결 서사. 저자 구성·데이터 공개 방식 협의", WARM),
]
y = 1.55
for code, title, desc, c in agenda:
    sidebar(s, 0.6, y, 1.05, c)
    runs(s, 0.85, y, 0.65, 0.4, [[(code, True, NAVY, 12)]])
    runs(s, 1.55, y, 2.85, 0.4, [[(title, True, NAVY, 11.5)]])
    runs(s, 4.5, y, 8.15, 0.95, [[(desc, False, INK, 9.8)]])
    y += 1.28
runs(s, 0.6, 6.7, 12.0, 0.3, [[("모든 실험은 커밋 단위로 기록되어 있어 어느 시점의 수치든 재현 가능합니다 — 회의에서 요청 주시면 현장에서 확인 가능합니다.", False, GREY, 10)]])
footnote(s, RULE)
notes(s, "우선 협의 대상은 A1(역할 분담). A3(병행 산출)는 부담 없는 수위로 제시 — 교체 아님 강조. A4 논문은 저자 구성이 민감할 수 있어 그들 페이스에 맞춤.")

out = "/Users/user/vibe/bistro-lstm/docs/한은미팅_적응학습경과_2026-09-04.pptx"
prs.save(out); print("saved:", out)
