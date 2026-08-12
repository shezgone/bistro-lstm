# -*- coding: utf-8 -*-
"""1p 리포트 — 동일 결합 규칙에서 Chronos-2f vs BISTRO 부품 교체 비교 (실무자 스타일)"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

F = "Apple SD Gothic Neo"
INK = RGBColor(0x26,0x26,0x26); GREY = RGBColor(0x6B,0x6B,0x6B)
LINE = RGBColor(0xC9,0xC9,0xC9); NAVY = RGBColor(0x24,0x34,0x47)
GREEN = RGBColor(0x00,0x8A,0x3E); BGGREY = RGBColor(0xF5,0xF5,0xF4)
WARM = RGBColor(0xB0,0x53,0x2F); WHITE = RGBColor(0xFF,0xFF,0xFF)

SC = "/private/tmp/claude-502/-Users-user-vibe-bistro-lstm/031ca4b9-003d-4cb2-afaa-06353432dc9b/scratchpad"
prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
s = prs.slides.add_slide(prs.slide_layouts[6])

def runs(x, y, w, h, lines, size=10, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp=2):
    tf = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.02); tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01); tf.margin_bottom = Inches(0.01)
    for i, rs in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp)
        for t, b, c, z in rs:
            r = p.add_run(); r.text = t
            r.font.name = F; r.font.size = Pt(z if z else size)
            r.font.bold = b; r.font.color.rgb = c if c else INK
            rp = r._r.get_or_add_rPr(); rp.append(rp.makeelement(qn('a:ea'), {'typeface': F}))
    return tf

def hline(x, y, w, color=LINE, pt=0.75):
    ln = s.shapes.add_connector(1, Inches(x), Inches(y), Inches(x+w), Inches(y))
    ln.line.color.rgb = color; ln.line.width = Pt(pt); return ln

def rect(x, y, w, h, fill=None, line=None):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False; return shp

# 헤더
runs(0.6, 0.32, 11.5, 0.3, [[("부품 교체 실험  |  네이버클라우드 AX Forward Lab · 2026. 8.", True, GREY, 10.5)]])
runs(0.6, 0.6, 12.1, 0.6, [[("같은 결합 규칙에 부품만 바꾸면 — Chronos-2f는 개선, BISTRO는 악화됩니다", True, NAVY, 17)]])
hline(0.6, 1.22, 12.13, NAVY, 1.2)

# 실험 설계 한 줄
runs(0.6, 1.38, 12.1, 0.35, [
    [("실험 설계: 검증된 주차별 결합에서 ", False, INK, 10.5),
     ("조기 슬롯의 부품 하나만 교체", True, INK, 10.5),
     ("하고 나머지는 전부 동일 — 조기(발표 19~14주 전) = (GBM + [부품])÷2, 이후 = XGBoost 단독. 부품 = Chronos-2f 또는 BISTRO(이식판).", False, INK, 10.5)],
])

# 좌: 차트
s.shapes.add_picture(f"{SC}/c2f_vs_bistro.png", Inches(0.5), Inches(1.95), width=Inches(7.3))
# 좌하: 원인 진단 표
runs(0.6, 5.35, 7.0, 0.35, [[("원인 진단 — 조기 구간 단독 성능이 갈림길입니다", True, NAVY, 11.5)]])
hline(0.6, 5.7, 7.1, LINE, 0.75)
tbl = [("조기 구간(-19~-14주) 단독 RMSE", "GBM 0.997", "Chronos-2f 0.999", "XGBoost 1.024", "BISTRO 1.277")]
runs(0.6, 5.82, 7.1, 0.4, [
    [("조기 구간 단독 RMSE   ", True, GREY, 9.5),
     ("GBM 0.997  ·  Chronos-2f 0.999  ·  XGBoost 1.024  ·  ", False, INK, 10),
     ("BISTRO 1.277", True, WARM, 10)],
    [("Chronos-2f는 이 구간에서 GBM과 대등(±0.002)하지만, BISTRO는 1.28배 열세 — 평균에 넣는 순간 조합 전체를 끌어내립니다.", False, INK, 9.5)],
], sp=3)

# 우: 해석
X, W = 8.15, 4.6
rect(X, 1.95, 0.045, 2.15, GREEN)
runs(X + 0.2, 1.95, W - 0.2, 2.2, [
    [("결합 규칙의 공이 아니라, 부품의 역량입니다", True, NAVY, 11.5)],
    [("주차별 결합(-1.4%)의 개선분은 “섞는 기술”이 아니라", False, INK, 10)],
    [("Chronos-2f가 지표 공백 구간에서 실제로 경쟁력이 있기", False, INK, 10)],
    [("때문입니다. 같은 자리에 BISTRO를 넣으면 동일 규칙이", False, INK, 10)],
    [("+3.2% 악화로 뒤집힙니다 (반등 분기는 0.85 → 1.01).", False, INK, 10)],
    [("", False, INK, 4)],
    [("오차 상관도 답을 바꾸지 못합니다 — BISTRO가 GBM과", False, INK, 10)],
    [("덜 겹치지만(0.904 vs 0.959), 다양성 이득이 정확도", False, INK, 10)],
    [("격차를 메우기엔 턱없이 부족합니다.", False, INK, 10)],
], sp=2)
rect(X, 4.35, 0.045, 1.5, NAVY)
runs(X + 0.2, 4.35, W - 0.2, 1.55, [
    [("시사점", True, NAVY, 11.5)],
    [("조기 슬롯 부품의 자격 요건은 “지표 공백 구간에서", False, INK, 10)],
    [("GBM과 대등한 단독 성능”입니다. 사전학습 계열은 이를", False, INK, 10)],
    [("충족하고, 소표본 직접학습 계열은 충족하지 못합니다 —", False, INK, 10)],
    [("7월 전수 검증의 결론이 부품 단위에서도 재확인됩니다.", False, INK, 10)],
], sp=2)
rect(X, 6.05, W, 0.78, BGGREY, None)
rect(X, 6.05, 0.045, 0.78, GREY)
runs(X + 0.2, 6.13, W - 0.4, 0.65, [
    [("유의: BISTRO는 당사 GDP 이식 재구현판 기준이며, 원 연구·타 과제", False, GREY, 8.8)],
    [("성능에 대한 판정이 아닙니다. 원본도 동일 규약 채점 시 비교 가능합니다.", False, GREY, 8.8)],
], sp=1.5)

# 각주
runs(0.6, 7.06, 11.5, 0.3, [
    [("주: 공통표본 26개 분기(BISTRO 예측이 없는 2018Q1~2019Q2 제외) — 32개 분기 리더보드 수치(XGBoost 0.750 등)와 직접 비교 불가. "
      "실시간 빈티지 · 속보치 · w[-19,-1] 평균 RMSE · schema v2. 낮을수록 정확.", False, GREY, 8)]])
runs(12.35, 7.06, 0.5, 0.3, [[("1", False, GREY, 9)]], align=PP_ALIGN.RIGHT)

out = "/Users/user/vibe/bistro-lstm/docs/부품비교_Chronos2f_vs_BISTRO_1p_2026-08-12.pptx"
prs.save(out); print("saved:", out)
