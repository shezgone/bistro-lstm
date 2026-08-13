# -*- coding: utf-8 -*-
"""1p 리포트 v2 — FM 부품 비교 (동일 프로토콜): Chronos-2f vs Moirai vs 직접학습 LSTM"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

F = "Apple SD Gothic Neo"
INK = RGBColor(0x26,0x26,0x26); GREY = RGBColor(0x6B,0x6B,0x6B)
LINE = RGBColor(0xC9,0xC9,0xC9); NAVY = RGBColor(0x24,0x34,0x47)
GREEN = RGBColor(0x00,0x8A,0x3E); BGGREY = RGBColor(0xF5,0xF5,0xF4)
WARM = RGBColor(0xB0,0x53,0x2F); WHITE = RGBColor(0xFF,0xFF,0xFF)

SC = "/private/tmp/claude-502/-Users-user-vibe-bistro-lstm/031ca4b9-003d-4cb2-afaa-06353432dc9b/scratchpad"
prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
s = prs.slides.add_slide(prs.slide_layouts[6])

def runs(x, y, w, h, lines, size=10, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp=2):
    tf = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.02); tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01); tf.margin_bottom = Inches(0.01)
    for i, rs in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp)
        for t, b, c, z in rs:
            r = p.add_run(); r.text = t
            r.font.name = F; r.font.size = Pt(z if z else size)
            r.font.bold = b; r.font.color.rgb = c if c else INK
            rp = r._r.get_or_add_rPr(); rp.append(rp.makeelement(qn('a:ea'), {'typeface': F}))
    return tf

def hline(x, y, w, color=LINE, pt=0.75):
    ln = s.shapes.add_connector(1, Inches(x), Inches(y), Inches(x+w), Inches(y))
    ln.line.color.rgb = color; ln.line.width = Pt(pt); return ln

def rect(x, y, w, h, fill=None, line=None):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False; return shp

# 헤더
runs(0.6, 0.32, 11.5, 0.3, [[("부품 교체 실험 — BISTRO 원본 실측판 (완전 동일 프로토콜)  |  네이버클라우드 AX Forward Lab · 2026. 8.", True, GREY, 10.5)]])
runs(0.6, 0.6, 12.1, 0.6, [[("BISTRO 원본을 직접 실측했습니다 — 조기 슬롯 자격을 통과한 부품은 여전히 Chronos-2f뿐", True, NAVY, 17)]])
hline(0.6, 1.22, 12.13, NAVY, 1.2)

# 실험 설계
runs(0.6, 1.38, 12.1, 0.6, [
    [("실험 설계: BISTRO 프레임워크의 세 암(Task-Specific 소형 모델 · Foundation Model · LLM ICL) 중 채점 가능한 부품들을 ", False, INK, 10.5),
     ("전부 동일 프로토콜", True, INK, 10.5),
     ("(DFM 스냅샷 + 공변량 10종 + 빠른신호 4종, 분기말 N_gdp 외삽, 32분기 실시간 그리드)로 실행하고, 검증된 주차별 결합의 조기 슬롯에 교체 장착.", False, INK, 10.5)],
    [("FM 부품 = Chronos-2f(120M, zero-shot) vs ", False, GREY, 9.5),
     ("BISTRO 원본", True, GREY, 9.5),
     ("(BIS WP 1337 공개 체크포인트 — Moirai-base 91M을 BIS 거시 4,925계열로 미세조정, zero-shot) · 참고 = 직접학습 LSTM · LLM ICL 암은 오염 문제로 채점 불성립(제외).", False, GREY, 9.5)],
])
# 좌: 차트
s.shapes.add_picture(f"{SC}/fm_parts.png", Inches(0.5), Inches(2.15), width=Inches(7.3))
# 좌하: 원본 실측 확보 경위
runs(0.6, 5.75, 7.1, 0.32, [[("이번 판의 핵심 — 프록시가 아니라 BISTRO 원본입니다", True, NAVY, 11)]])
hline(0.6, 6.07, 7.1, LINE, 0.75)
runs(0.6, 6.16, 7.1, 0.9, [
    [("· BISTRO 가중치를 BIS 공개 레포(bis-med-it/bistro, Apache 2.0)에서 확보해 직접 실측 — 단독 ", False, INK, 9.3),
     ("0.871", True, INK, 9.3), (" (기반 Moirai-small 0.873과 동급, Chronos-2f 0.808에 미달).", False, INK, 9.3)],
    [("· 거시 미세조정의 효과는 소폭(조기 구간 1.021→1.005) — 미세조정으로도 아래의 구조 차이는 바뀌지 않았습니다.", False, INK, 9.3)],
    [("· 참고: 직접학습 LSTM은 동일 입력에도 1.019 (조기 1.286) — 소표본 직접학습 한계 별도 확인.", False, INK, 9.3)],
], sp=2)

# 우측
X, W = 8.15, 4.6
rect(X, 2.15, 0.045, 1.5, GREEN)
runs(X + 0.2, 2.15, W - 0.2, 1.55, [
    [("슬롯 자격 요건은 여전히 하나입니다", True, NAVY, 11.5)],
    [("“조기 구간(지표 공백기)에서 GBM과 대등한 단독 성능”", False, INK, 9.5)],
    [("— Chronos-2f 0.924 ≈ GBM 0.924 (통과) / BISTRO 1.005 ·", False, INK, 9.5)],
    [("LSTM 1.286 (미달). 슬롯 결과: C2f 0.740(개선) / BISTRO", False, INK, 9.5)],
    [("0.750(기준선과 동률 — 효과 없음) / 듀오 0.744 — C2f 단독을", False, INK, 9.5)],
    [("못 넘습니다 (오차 상관 0.926 — 정보 중복).", False, INK, 9.5)],
], sp=2)
rect(X, 3.8, 0.045, 2.5, NAVY)
runs(X + 0.2, 3.8, W - 0.2, 2.55, [
    [("무엇이 다른가 — 절제 실험이 답합니다 (힌트 흡수력)", True, NAVY, 11.5)],
    [("둘 다 공변량을 입력받지만, ", False, INK, 9)],
    [("같은 힌트에서 뽑아내는 이득", True, INK, 9), ("이 실측으로 다릅니다:", False, INK, 9)],
    [("· 빠른신호(주가·환율)를 얹었을 때 — ", False, INK, 9),
     ("C2 -3.4%", True, GREEN, 9), (" (0.836→0.808)", False, GREY, 8.5)],
    [("  vs ", False, INK, 9), ("BISTRO 무이득", True, WARM, 9), (" (0.870→0.871, +0.1%)", False, INK, 9)],
    [("· 공식지표 10종: BISTRO -0.5% (0.875→0.870) — 미미", False, INK, 9)],
    [("원인 — 사전학습에서 채점 범위가 다릅니다:", False, INK, 9)],
    [("· Moirai/BISTRO:  (공변량₁·₂, 타깃)과거 → ", False, INK, 9),
     ("(공변량₁·₂, 타깃)미래 전부 채점", True, INK, 9)],
    [("  = 공변량도 함께 예측할 대상(이웃). 거시 미세조정으로도", False, INK, 9)],
    [("  이 채점 구조는 그대로 → 힌트 활용 회로가 약하게 유지", False, INK, 9)],
    [("· Chronos-2:  (공변량₁·₂, 타깃)과거 → ", False, INK, 9),
     ("타깃 미래만 채점", True, INK, 9)],
    [("  = 공변량은 힌트 — 힌트를 안 보면 손실을 줄일 수 없게", False, INK, 9)],
    [("  의존관계가 설계된 사례(합성 포함)로 훈련", False, INK, 9)],
    [("· 체급은 비슷(91M vs 120M) — 격차 원인은 체급이 아니라 훈련 방식", False, GREY, 8.5)],
], sp=1.5)
rect(X, 6.42, W, 0.72, BGGREY, None)
rect(X, 6.42, 0.045, 0.72, GREY)
runs(X + 0.2, 6.47, W - 0.4, 0.65, [
    [("유의: BISTRO = BIS WP 1337 (Koyuncu·Kwon·Lombardi·Perez-Cruz·Shin)의", False, GREY, 8.5)],
    [("공개 원본 체크포인트를 논문 표준 사용법(zero-shot)으로 적용한 실측입니다.", False, GREY, 8.5)],
    [("한국 GDP 과제용 추가 미세조정은 미실시. LLM ICL 암은 오염으로 제외.", False, GREY, 8.5)],
], sp=1.5)

# 각주
runs(0.6, 7.22, 11.5, 0.25, [
    [("주: 실시간 빈티지 · 속보치 · w[-19,-1] 평균 RMSE · 2018Q1~2025Q4(32개 분기) · schema v2 · 낮을수록 정확. "
      "슬롯 교체 = early(발표 19~14주 전)=(GBM+부품)÷2, 이후 XGBoost. 개선폭은 모두 통계적 유의성 미달(DM p>0.18) — 비열등+개선 방향 수위.", False, GREY, 7.5)]])
runs(12.35, 7.2, 0.5, 0.25, [[("1", False, GREY, 9)]], align=PP_ALIGN.RIGHT)

out = "/Users/user/vibe/bistro-lstm/docs/부품비교_FM동일프로토콜_1p_2026-08-13.pptx"
prs.save(out); print("saved:", out)
