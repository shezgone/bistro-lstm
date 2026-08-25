# -*- coding: utf-8 -*-
"""회의 대응 1p — 주차별 빈티지 정보 흡수: BISTRO 평탄성 진단에 대한 검증과 제안"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

F = "Apple SD Gothic Neo"
INK = RGBColor(0x26,0x26,0x26); GREY = RGBColor(0x6B,0x6B,0x6B)
LINE = RGBColor(0xC9,0xC9,0xC9); NAVY = RGBColor(0x24,0x34,0x47)
GREEN = RGBColor(0x00,0x8A,0x3E); BGGREY = RGBColor(0xF5,0xF5,0xF4)
WARM = RGBColor(0xB0,0x53,0x2F)

SC = "/private/tmp/claude-502/-Users-user-vibe-bistro-lstm/031ca4b9-003d-4cb2-afaa-06353432dc9b/scratchpad"
prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
s = prs.slides.add_slide(prs.slide_layouts[6])

def runs(x, y, w, h, lines, size=10, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp=2):
    tf = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.02); tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01); tf.margin_bottom = Inches(0.01)
    for i, rs in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp)
        for t, b, c, z in rs:
            r = p.add_run(); r.text = t
            r.font.name = F; r.font.size = Pt(z if z else size)
            r.font.bold = b; r.font.color.rgb = c if c else INK
            rp = r._r.get_or_add_rPr(); rp.append(rp.makeelement(qn('a:ea'), {'typeface': F}))
    return tf

def hline(x, y, w, color=LINE, pt=0.75):
    ln = s.shapes.add_connector(1, Inches(x), Inches(y), Inches(x+w), Inches(y))
    ln.line.color.rgb = color; ln.line.width = Pt(pt); return ln

def rect(x, y, w, h, fill=None, line=None):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False; return shp

# 헤더
runs(0.6, 0.32, 11.5, 0.3, [[("협업회의 대응 — 귀 리포트(8/25) 그림 1 관련  |  네이버클라우드 AX Forward Lab · 2026. 8.", True, GREY, 10.5)]])
runs(0.6, 0.6, 12.1, 0.6, [[("주차별 평탄성은 모델 고유 특성 — 같은 조건에서 Chronos-2는 빈티지 갱신을 흡수", True, NAVY, 17)]])
hline(0.6, 1.22, 12.13, NAVY, 1.2)

# 검증 설계
runs(0.6, 1.38, 12.1, 0.6, [
    [("검증 설계  ", True, NAVY, 10.5),
     ("귀 리포트의 진단(BISTRO 오차가 전망주차 경과에도 감소하지 않음)을 당사 규약에서 재검증 — 동일 입력(DFM 스냅샷 + 공변량 10종 + 일별신호 4종)·동일 그리드(32분기×19주)에서 "
      "BISTRO와 Chronos-2의 주차별 RMSE 곡선을 비교. 분기말 월 관측 이후(w≥-7)는 TSFM이 패널 값을 그대로 쓰는 구간이라 자체 예측 구간(w=-19~-8)만 비교.", False, INK, 10.5)],
])
# 좌: 차트
s.shapes.add_picture(f"{SC}/weekly_absorb.png", Inches(0.5), Inches(2.1), width=Inches(7.5))
# 좌하: 흡수율 표 요약
runs(0.6, 6.12, 7.3, 0.3, [[("정보 흡수율 (초반 6주 → 다음 6주 RMSE 감소율)", True, NAVY, 10.5)]])
hline(0.6, 6.42, 7.3, LINE, 0.75)
runs(0.6, 6.5, 7.3, 0.6, [
    [("XGBoost -24.8%  ·  DFM -18.6%  ·  ", False, INK, 9.5),
     ("Chronos-2f -9.3%", True, GREEN, 9.5),
     ("  ·  Moirai-small -7.0%  ·  ", False, INK, 9.5),
     ("BISTRO -4.7%", True, WARM, 9.5)],
    [("Chronos-2 기본형(일별신호 제외)도 -11.0% — 흡수는 일별신호가 아니라 모델의 성질", False, GREY, 9)],
], sp=2)

# 우측
X, W = 8.35, 4.4
rect(X, 2.1, 0.045, 1.45, WARM)
runs(X + 0.2, 2.1, W - 0.2, 1.5, [
    [("진단 재현 — 평탄성은 입력 문제가 아님", True, NAVY, 11.5)],
    [("· 일별신호까지 준 동일 조건에서도 BISTRO는 -19~-9주", False, INK, 9.5)],
    [("  구간 0.96~1.06에서 평탄 (흡수율 -4.7%)", False, INK, 9.5)],
    [("· 기반 Moirai-small도 -7.0% — 계열 공통 특성", False, INK, 9.5)],
    [("· 예측 지평 단축의 기계적 이점을 받고도 평탄 →", False, INK, 9.5)],
    [("  귀 리포트의 '정보 미흡수' 진단과 일치", False, INK, 9.5)],
], sp=2)
rect(X, 3.7, 0.045, 1.5, GREEN)
runs(X + 0.2, 3.7, W - 0.2, 1.55, [
    [("같은 조건의 Chronos-2는 우하향", True, NAVY, 11.5)],
    [("· 흡수율 -9.3% (BISTRO의 약 2배), 레벨도 전 구간 우위", False, INK, 9.5)],
    [("· 원인은 사전학습 채점 범위의 차이 — 별첨(부품 비교)의", False, INK, 9.5)],
    [("  절제 실험과 정합: 같은 일별신호에서 C2 -3.4% vs", False, INK, 9.5)],
    [("  BISTRO 무이득", False, INK, 9.5)],
    [("· 단서: C2도 XGBoost(-24.8%)에는 미달 — TSFM 전반이", False, INK, 9.5)],
    [("  새 빈티지를 완전히 소화하지는 못함", False, INK, 9.5)],
], sp=2)
rect(X, 5.35, 0.045, 1.45, NAVY)
runs(X + 0.2, 5.35, W - 0.2, 1.5, [
    [("회의 제안", True, NAVY, 11.5)],
    [("· 적응학습(adapter/LoRA) 검토에 축 하나 추가 —", False, INK, 9.5)],
    [("  ", False, INK, 9.5), ("출발 체크포인트 선택", True, INK, 9.5),
     (": 흡수력이 확인된 계열(Chronos-2)의", False, INK, 9.5)],
    [("  zero-shot·적응 병행 비교를 제안", False, INK, 9.5)],
    [("· 단기 정확도 기준축은 XGBoost 유지에 동의", False, INK, 9.5)],
    [("· 당사 기여 가능: 흡수력 진단 하네스 · 일별신호 파이프라인 ·", False, INK, 9.5)],
    [("  절제 실험 프레임 · 공정 비교 규약", False, INK, 9.5)],
], sp=2)
rect(X, 6.9, W, 0.42, BGGREY, None)
rect(X, 6.9, 0.045, 0.42, GREY)
runs(X + 0.2, 6.94, W - 0.4, 0.35, [
    [("유의: 사전학습 기간과 평가기간의 중복 가능성은 귀 리포트와 동일하게", False, GREY, 8)],
    [("회고적 진단으로 해석. 모든 개선폭은 통계적 유의성 미달 수위 유지.", False, GREY, 8)],
], sp=1.2)

# 각주 (우측 유의 박스와 겹치지 않게 좌측 폭으로 제한)
runs(0.6, 7.1, 7.5, 0.35, [
    [("주: 실시간 빈티지 · 속보치 · 32분기 주차별 RMSE(분기 평균) · schema v2 · 낮을수록 정확. "
      "Chronos-2f = Chronos-2(120M)+일별신호 · BISTRO = BIS WP 1337 공개 체크포인트(91M), 모두 zero-shot.", False, GREY, 7.5)]])
runs(12.85, 7.15, 0.4, 0.25, [[("1", False, GREY, 9)]], align=PP_ALIGN.RIGHT)

out = "/Users/user/vibe/bistro-lstm/docs/회의대응_빈티지흡수_1p_2026-08-26.pptx"
prs.save(out); print("saved:", out)
